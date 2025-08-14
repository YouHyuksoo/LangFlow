"""
스트리밍 청크 처리 시스템
- 메모리 효율적인 대용량 파일 처리
- 실시간 진행 상황 업데이트
- 백프레셔 제어
"""
import asyncio
import time
from typing import AsyncGenerator, Dict, Any, List, Optional, Callable
from dataclasses import dataclass
from pathlib import Path
import logging

from ..core.config import settings


@dataclass
class StreamingChunk:
    """스트리밍 청크 데이터"""
    content: str
    index: int
    total_chunks: int
    metadata: Dict[str, Any]
    timestamp: float


@dataclass 
class ProcessingProgress:
    """처리 진행 상황"""
    processed_chunks: int
    total_chunks: int
    processing_rate: float  # chunks/second
    estimated_remaining_time: float  # seconds
    current_stage: str
    errors: List[str]


class StreamingChunkProcessor:
    """스트리밍 청크 처리기"""
    
    def __init__(self, progress_callback: Optional[Callable] = None):
        self.logger = logging.getLogger(__name__)
        self.progress_callback = progress_callback
        self.buffer_size = settings.CHUNK_STREAM_BUFFER_SIZE
        
        # 진행 상황 추적
        self.start_time = 0
        self.processed_count = 0
        self.error_count = 0
        self.processing_rates = []  # 최근 처리 속도 기록
        
    async def stream_process_file(
        self, 
        file_path: str,
        chunk_size: int = None,
        overlap_size: int = None
    ) -> AsyncGenerator[StreamingChunk, None]:
        """파일을 스트리밍 방식으로 청크 처리"""
        chunk_size = chunk_size or settings.DEFAULT_CHUNK_SIZE
        overlap_size = overlap_size or settings.DEFAULT_CHUNK_OVERLAP
        
        self.logger.info(f"스트리밍 파일 처리 시작: {file_path}")
        self.start_time = time.time()
        
        try:
            # 파일 크기 확인
            file_size = Path(file_path).stat().st_size
            estimated_chunks = file_size // (chunk_size - overlap_size)
            
            self.logger.info(f"예상 청크 수: {estimated_chunks}개 (파일 크기: {file_size / 1024 / 1024:.1f}MB)")
            
            # 파일 확장자에 따른 처리 방식 결정
            file_extension = Path(file_path).suffix.lower()
            
            if file_extension == '.pdf':
                async for chunk in self._stream_pdf_chunks(file_path, chunk_size, overlap_size):
                    yield chunk
            elif file_extension in ['.docx', '.pptx', '.xlsx']:
                async for chunk in self._stream_office_chunks(file_path, chunk_size, overlap_size):
                    yield chunk
            else:
                async for chunk in self._stream_text_chunks(file_path, chunk_size, overlap_size):
                    yield chunk
                    
        except Exception as e:
            self.logger.error(f"스트리밍 처리 중 오류: {e}")
            raise
    
    async def _stream_pdf_chunks(
        self, 
        file_path: str, 
        chunk_size: int, 
        overlap_size: int
    ) -> AsyncGenerator[StreamingChunk, None]:
        """PDF 파일 스트리밍 청크 처리"""
        try:
            # PyPDF를 사용한 페이지별 스트리밍 처리
            from pypdf import PdfReader
            
            pdf_reader = PdfReader(file_path)
            total_pages = len(pdf_reader.pages)
            
            self.logger.info(f"PDF 스트리밍 처리 - 총 {total_pages}페이지")
            
            current_text = ""
            chunk_index = 0
            overlap_text = ""
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    # 진행 상황 업데이트
                    await self._update_progress(
                        processed=page_num,
                        total=total_pages,
                        stage=f"PDF 페이지 {page_num + 1}/{total_pages} 처리 중"
                    )
                    
                    # 페이지 텍스트 추출
                    page_text = page.extract_text()
                    if not page_text.strip():
                        continue
                    
                    # 오버랩과 함께 텍스트 결합
                    current_text = overlap_text + page_text
                    
                    # 청크 크기에 맞게 분할
                    while len(current_text) >= chunk_size:
                        chunk_text = current_text[:chunk_size]
                        
                        # 청크 생성 및 전송
                        chunk = StreamingChunk(
                            content=chunk_text,
                            index=chunk_index,
                            total_chunks=-1,  # 스트리밍이므로 미지
                            metadata={
                                "source_page": page_num + 1,
                                "chunk_type": "pdf_stream"
                            },
                            timestamp=time.time()
                        )
                        
                        yield chunk
                        chunk_index += 1
                        self.processed_count += 1
                        
                        # 오버랩 설정
                        overlap_text = current_text[chunk_size - overlap_size:]
                        current_text = current_text[chunk_size - overlap_size:]
                        
                        # 백프레셔 제어
                        if chunk_index % self.buffer_size == 0:
                            await asyncio.sleep(0.01)
                    
                    # 나머지 텍스트를 다음 페이지로 이월
                    overlap_text = current_text
                    
                except Exception as page_error:
                    self.logger.warning(f"페이지 {page_num + 1} 처리 실패: {page_error}")
                    self.error_count += 1
                    continue
            
            # 마지막 남은 텍스트 처리
            if current_text.strip():
                chunk = StreamingChunk(
                    content=current_text,
                    index=chunk_index,
                    total_chunks=chunk_index + 1,
                    metadata={"chunk_type": "pdf_stream_final"},
                    timestamp=time.time()
                )
                yield chunk
                self.processed_count += 1
                
        except Exception as e:
            self.logger.error(f"PDF 스트리밍 처리 실패: {e}")
            raise
    
    async def _stream_office_chunks(
        self, 
        file_path: str, 
        chunk_size: int, 
        overlap_size: int
    ) -> AsyncGenerator[StreamingChunk, None]:
        """Office 파일 스트리밍 청크 처리"""
        file_extension = Path(file_path).suffix.lower()
        
        try:
            if file_extension == '.docx':
                async for chunk in self._stream_docx_chunks(file_path, chunk_size, overlap_size):
                    yield chunk
            elif file_extension == '.pptx':
                async for chunk in self._stream_pptx_chunks(file_path, chunk_size, overlap_size):
                    yield chunk
            elif file_extension == '.xlsx':
                async for chunk in self._stream_xlsx_chunks(file_path, chunk_size, overlap_size):
                    yield chunk
                    
        except Exception as e:
            self.logger.error(f"Office 파일 스트리밍 처리 실패: {e}")
            raise
    
    async def _stream_docx_chunks(
        self, 
        file_path: str, 
        chunk_size: int, 
        overlap_size: int
    ) -> AsyncGenerator[StreamingChunk, None]:
        """DOCX 파일 스트리밍 처리"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            total_paragraphs = len(doc.paragraphs)
            
            current_text = ""
            chunk_index = 0
            overlap_text = ""
            
            for para_idx, paragraph in enumerate(doc.paragraphs):
                await self._update_progress(
                    processed=para_idx,
                    total=total_paragraphs,
                    stage=f"DOCX 단락 {para_idx + 1}/{total_paragraphs} 처리 중"
                )
                
                if not paragraph.text.strip():
                    continue
                
                current_text = overlap_text + paragraph.text + "\n"
                
                # 청크 분할
                while len(current_text) >= chunk_size:
                    chunk_text = current_text[:chunk_size]
                    
                    chunk = StreamingChunk(
                        content=chunk_text,
                        index=chunk_index,
                        total_chunks=-1,
                        metadata={
                            "source_paragraph": para_idx + 1,
                            "chunk_type": "docx_stream"
                        },
                        timestamp=time.time()
                    )
                    
                    yield chunk
                    chunk_index += 1
                    self.processed_count += 1
                    
                    overlap_text = current_text[chunk_size - overlap_size:]
                    current_text = current_text[chunk_size - overlap_size:]
                    
                    if chunk_index % self.buffer_size == 0:
                        await asyncio.sleep(0.01)
                
                overlap_text = current_text
            
            # 마지막 청크
            if current_text.strip():
                chunk = StreamingChunk(
                    content=current_text,
                    index=chunk_index,
                    total_chunks=chunk_index + 1,
                    metadata={"chunk_type": "docx_stream_final"},
                    timestamp=time.time()
                )
                yield chunk
                self.processed_count += 1
                
        except Exception as e:
            self.logger.error(f"DOCX 스트리밍 처리 실패: {e}")
            raise
    
    async def _stream_pptx_chunks(
        self, 
        file_path: str, 
        chunk_size: int, 
        overlap_size: int
    ) -> AsyncGenerator[StreamingChunk, None]:
        """PPTX 파일 스트리밍 처리"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            
            chunk_index = 0
            
            for slide_idx, slide in enumerate(prs.slides):
                await self._update_progress(
                    processed=slide_idx,
                    total=total_slides,
                    stage=f"PPTX 슬라이드 {slide_idx + 1}/{total_slides} 처리 중"
                )
                
                slide_text = f"[슬라이드 {slide_idx + 1}]\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if len(slide_text.strip()) > 10:  # 의미있는 내용이 있는 경우
                    chunk = StreamingChunk(
                        content=slide_text,
                        index=chunk_index,
                        total_chunks=-1,
                        metadata={
                            "source_slide": slide_idx + 1,
                            "chunk_type": "pptx_stream"
                        },
                        timestamp=time.time()
                    )
                    
                    yield chunk
                    chunk_index += 1
                    self.processed_count += 1
                    
                    if chunk_index % self.buffer_size == 0:
                        await asyncio.sleep(0.01)
                        
        except Exception as e:
            self.logger.error(f"PPTX 스트리밍 처리 실패: {e}")
            raise
    
    async def _stream_xlsx_chunks(
        self, 
        file_path: str, 
        chunk_size: int, 
        overlap_size: int
    ) -> AsyncGenerator[StreamingChunk, None]:
        """XLSX 파일 스트리밍 처리"""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(file_path, data_only=True)
            total_sheets = len(wb.sheetnames)
            
            chunk_index = 0
            
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                await self._update_progress(
                    processed=sheet_idx,
                    total=total_sheets,
                    stage=f"XLSX 시트 {sheet_idx + 1}/{total_sheets} 처리 중"
                )
                
                sheet = wb[sheet_name]
                sheet_text = f"[시트: {sheet_name}]\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None:
                            row_text.append(str(cell_value))
                    
                    if row_text:
                        sheet_text += " | ".join(row_text) + "\n"
                        
                        # 청크 크기 확인
                        if len(sheet_text) >= chunk_size:
                            chunk = StreamingChunk(
                                content=sheet_text,
                                index=chunk_index,
                                total_chunks=-1,
                                metadata={
                                    "source_sheet": sheet_name,
                                    "chunk_type": "xlsx_stream"
                                },
                                timestamp=time.time()
                            )
                            
                            yield chunk
                            chunk_index += 1
                            self.processed_count += 1
                            
                            sheet_text = f"[시트: {sheet_name}]\n"  # 리셋
                            
                            if chunk_index % self.buffer_size == 0:
                                await asyncio.sleep(0.01)
                
                # 시트별 마지막 청크
                if len(sheet_text.strip()) > len(f"[시트: {sheet_name}]"):
                    chunk = StreamingChunk(
                        content=sheet_text,
                        index=chunk_index,
                        total_chunks=-1,
                        metadata={
                            "source_sheet": sheet_name,
                            "chunk_type": "xlsx_stream"
                        },
                        timestamp=time.time()
                    )
                    
                    yield chunk
                    chunk_index += 1
                    self.processed_count += 1
                    
        except Exception as e:
            self.logger.error(f"XLSX 스트리밍 처리 실패: {e}")
            raise
    
    async def _stream_text_chunks(
        self, 
        file_path: str, 
        chunk_size: int, 
        overlap_size: int
    ) -> AsyncGenerator[StreamingChunk, None]:
        """일반 텍스트 파일 스트리밍 처리"""
        try:
            chunk_index = 0
            current_text = ""
            overlap_text = ""
            
            # 파일을 스트리밍으로 읽기
            with open(file_path, 'r', encoding='utf-8') as file:
                while True:
                    # 한 번에 읽을 크기 (메모리 효율성)
                    read_size = chunk_size
                    data = file.read(read_size)
                    
                    if not data:
                        break
                    
                    current_text = overlap_text + data
                    
                    # 청크 분할
                    while len(current_text) >= chunk_size:
                        chunk_text = current_text[:chunk_size]
                        
                        chunk = StreamingChunk(
                            content=chunk_text,
                            index=chunk_index,
                            total_chunks=-1,
                            metadata={"chunk_type": "text_stream"},
                            timestamp=time.time()
                        )
                        
                        yield chunk
                        chunk_index += 1
                        self.processed_count += 1
                        
                        overlap_text = current_text[chunk_size - overlap_size:]
                        current_text = current_text[chunk_size - overlap_size:]
                        
                        if chunk_index % self.buffer_size == 0:
                            await asyncio.sleep(0.01)
                    
                    overlap_text = current_text
            
            # 마지막 청크
            if current_text.strip():
                chunk = StreamingChunk(
                    content=current_text,
                    index=chunk_index,
                    total_chunks=chunk_index + 1,
                    metadata={"chunk_type": "text_stream_final"},
                    timestamp=time.time()
                )
                yield chunk
                self.processed_count += 1
                
        except Exception as e:
            self.logger.error(f"텍스트 스트리밍 처리 실패: {e}")
            raise
    
    async def _update_progress(self, processed: int, total: int, stage: str):
        """진행 상황 업데이트"""
        if not self.progress_callback:
            return
        
        elapsed_time = time.time() - self.start_time
        if elapsed_time > 0:
            processing_rate = self.processed_count / elapsed_time
            self.processing_rates.append(processing_rate)
            
            # 최근 10개 측정값으로 평균 계산
            if len(self.processing_rates) > 10:
                self.processing_rates = self.processing_rates[-10:]
            
            avg_rate = sum(self.processing_rates) / len(self.processing_rates)
            remaining_chunks = total - processed
            estimated_remaining_time = remaining_chunks / avg_rate if avg_rate > 0 else 0
            
            progress = ProcessingProgress(
                processed_chunks=self.processed_count,
                total_chunks=total,
                processing_rate=avg_rate,
                estimated_remaining_time=estimated_remaining_time,
                current_stage=stage,
                errors=[]
            )
            
            try:
                await self.progress_callback(progress)
            except Exception as e:
                self.logger.warning(f"진행 상황 콜백 오류: {e}")
    
    def get_statistics(self) -> Dict[str, Any]:
        """처리 통계 반환"""
        elapsed_time = time.time() - self.start_time
        return {
            "processed_chunks": self.processed_count,
            "error_count": self.error_count,
            "elapsed_time": elapsed_time,
            "processing_rate": self.processed_count / elapsed_time if elapsed_time > 0 else 0,
            "memory_efficient": True,
            "streaming_mode": True
        }