
import os
import uuid
import aiofiles
import time
import logging
import asyncio
from typing import List, Optional, Dict, Any
from fastapi import UploadFile, HTTPException
from datetime import datetime

from ..models.schemas import FileUploadResponse, FileInfo, FileProcessingOptions
from ..models.vector_models import FileMetadata, FileMetadataService
from ..models.schemas import FileStatus
from ..core.config import settings
from .category_service import CategoryService
from .preprocessing_service import PreprocessingService
from .vector_service import VectorService
from .settings_service import settings_service

# SSE 이벤트 전송용
try:
    from ..api.sse import get_sse_manager
    SSE_AVAILABLE = True
except ImportError:
    SSE_AVAILABLE = False
    print("SSE 모듈을 찾을 수 없습니다.")

class FileService:
    """파일 관리 및 벡터화 파이프라인 오케스트레이션을 담당합니다."""
    def __init__(self):
        self.upload_dir = settings.UPLOAD_DIR
        self.logger = logging.getLogger(__name__)
        self._load_allowed_extensions()
        self.category_service = CategoryService()
        
        # Refactored services
        self.preprocessing_service = PreprocessingService()
        self.vector_service = VectorService()

        # SQLite 기반 파일 메타데이터 서비스
        self.file_metadata_service = FileMetadataService()
        self._ensure_data_dir()

    # --- 분리된 전처리 및 벡터화 파이프라인 ---
    async def start_preprocessing(self, file_id: str, method: str = None):
        """파일 전처리를 시작합니다."""
        preprocessing_start_time = time.time()
        self.logger.info(f"🔄 === 전처리 시작: {file_id} ===")

        try:
            # method가 지정되지 않았으면 설정에서 가져오기
            if method is None:
                from .settings_service import settings_service
                system_settings = settings_service.get_section_settings("system")
                method = system_settings.get("preprocessing_method", "basic")
                self.logger.info(f"설정에서 가져온 전처리 방식: {method}")
            
            file_info = await self.get_file_info(file_id)
            if not file_info or not file_info.file_path or not os.path.exists(file_info.file_path):
                self.logger.error(f"파일 정보를 찾을 수 없거나 파일이 존재하지 않습니다: {file_id}")
                await self._update_file_status(file_id, FileStatus.FAILED, error="File not found or path is invalid.")
                return {"success": False, "error": "File not found"}

            if file_info.status != FileStatus.UPLOADED:
                self.logger.warning(f"파일이 전처리 가능한 상태가 아닙니다: {file_id}, 현재 상태: {file_info.status}")
                return {"success": False, "error": "File is not in uploadable state"}

            await self._update_file_status(file_id, FileStatus.PREPROCESSING)

            # 전처리 실행
            self.logger.info(f"텍스트 추출 및 전처리 시작 (방법: {method})...")
            text_content = await self.preprocessing_service.process_file(file_info.file_path, method)
            
            # 전처리 결과 저장
            await self._save_preprocessed_content(file_id, text_content, method)
            
            await self._update_file_status(file_id, FileStatus.PREPROCESSED)
            
            elapsed = time.time() - preprocessing_start_time
            self.logger.info(f"✅ 전처리 완료. 추출된 텍스트 길이: {len(text_content)} (소요 시간: {elapsed:.2f}초)")
            
            return {
                "success": True,
                "text_length": len(text_content),
                "processing_time": elapsed,
                "method": method
            }

        except Exception as e:
            self.logger.error(f"💥 전처리 오류: {e}", exc_info=True)
            await self._update_file_status(file_id, FileStatus.FAILED, error=str(e))
            return {"success": False, "error": str(e)}

    async def start_vectorization(self, file_id: str):
        """전처리된 파일의 벡터화를 시작합니다."""
        vectorization_start_time = time.time()
        self.logger.info(f"🚀 === 벡터화 시작: {file_id} ===")

        try:
            file_info = await self.get_file_info(file_id)
            if not file_info:
                self.logger.error(f"파일 정보를 찾을 수 없습니다: {file_id}")
                return {"success": False, "error": "File not found"}

            # PREPROCESSING 상태나 FAILED 상태의 파일도 재처리 가능하도록 허용
            if file_info.status not in [FileStatus.UPLOADED, FileStatus.PREPROCESSED, FileStatus.VECTORIZING, FileStatus.PREPROCESSING, FileStatus.FAILED]:
                self.logger.warning(f"파일이 벡터화 가능한 상태가 아닙니다: {file_id}, 현재 상태: {file_info.status}")
                return {"success": False, "error": "File is not ready for vectorization"}
            
            # PREPROCESSING 또는 FAILED 상태인 경우 강제 재처리 시작
            if file_info.status in [FileStatus.PREPROCESSING, FileStatus.FAILED]:
                self.logger.info(f"🔄 상태가 {file_info.status}인 파일을 강제 재처리합니다: {file_id}")
                # 상태를 UPLOADED로 재설정하여 처음부터 다시 시작
                await self._update_file_status(file_id, FileStatus.UPLOADED)

            # UPLOADED 상태라면 전처리부터 시작
            if file_info.status == FileStatus.UPLOADED:
                self.logger.info(f"UPLOADED 상태 파일 전처리 시작: {file_id}")
                # 기본 설정에서 전처리 방법 가져오기 (설정에 따라 basic, docling, unstructured 중 선택)
                from .settings_service import settings_service
                system_settings = settings_service.get_section_settings("system")
                preprocessing_method = system_settings.get("preprocessing_method", "basic")
                self.logger.info(f"설정된 전처리 방식: {preprocessing_method}")
                preprocess_result = await self.start_preprocessing(file_id, preprocessing_method)
                if not preprocess_result.get("success"):
                    self.logger.error(f"전처리 실패: {file_id}")
                    return {"success": False, "error": "Preprocessing failed"}
                
                # 전처리 완료 후 파일 정보 다시 로드
                file_info = await self.get_file_info(file_id)

            await self._update_file_status(file_id, FileStatus.VECTORIZING)

            # SSE 이벤트 전송 (벡터화 시작)
            if SSE_AVAILABLE:
                try:
                    sse_manager = get_sse_manager()
                    await sse_manager.broadcast("vectorization_update", {
                        "file_id": file_id,
                        "filename": file_info.filename,
                        "status": "started",
                        "vectorized": False
                    })
                    self.logger.info(f"📡 SSE 벡터화 시작 이벤트 전송: {file_id}")
                except Exception as sse_error:
                    self.logger.warning(f"SSE 시작 이벤트 전송 실패: {sse_error}")

            # 전처리된 텍스트 로드
            text_content = await self._load_preprocessed_content(file_id)
            if not text_content:
                self.logger.error(f"전처리된 텍스트를 찾을 수 없습니다: {file_id}")
                await self._update_file_status(file_id, FileStatus.FAILED, error="Preprocessed content not found")
                return {"success": False, "error": "Preprocessed content not found"}

            # 벡터화 실행
            self.logger.info(f"청킹 및 임베딩 시작...")
            vector_metadata = { 
                "filename": file_info.filename, 
                "category_id": file_info.category_id,
                "category_name": file_info.category_name,
                "preprocessing_method": file_info.preprocessing_method
            }
            result = await self.vector_service.chunk_and_embed_text(file_id, text_content, vector_metadata)

            if result.get("success"):
                self.logger.info(f"🔄 메타데이터 업데이트 시작 - 파일 상태를 COMPLETED로 변경")
                await self._update_file_status(file_id, FileStatus.COMPLETED, chunks_count=result.get('chunks_count'))
                
                # vectorized 필드도 별도로 업데이트
                self.logger.info(f"🔄 vectorized 상태 업데이트 시작")
                await self.update_file_vectorization_status(
                    file_id=file_id,
                    vectorized=True,
                    error_message=None,
                    chunk_count=result.get('chunks_count', 0)
                )
                
                elapsed = time.time() - vectorization_start_time
                self.logger.info(f"📋 메타데이터 업데이트 완료 - 청크 수: {result.get('chunks_count', 0)}")
                self.logger.info(f"✅ 벡터화 성공. 청크 수: {result.get('chunks_count', 0)} (소요 시간: {elapsed:.2f}초)")
                
                # SSE 이벤트 전송 (벡터화 완료)
                if SSE_AVAILABLE:
                    try:
                        sse_manager = get_sse_manager()
                        await sse_manager.broadcast("vectorization_update", {
                            "file_id": file_id,
                            "filename": file_info.filename,
                            "status": "completed",
                            "vectorized": True,
                            "chunks_count": result.get('chunks_count', 0),
                            "processing_time": elapsed
                        })
                        self.logger.info(f"📡 SSE 벡터화 완료 이벤트 전송: {file_id}")
                    except Exception as sse_error:
                        self.logger.warning(f"SSE 이벤트 전송 실패: {sse_error}")
                
                return {
                    "success": True,
                    "chunks_count": result.get('chunks_count', 0),
                    "processing_time": elapsed
                }
            else:
                self.logger.error(f"❌ 벡터화 실패: {result.get('error')}")
                await self._update_file_status(file_id, FileStatus.FAILED, error=result.get('error'))
                
                # SSE 이벤트 전송 (벡터화 실패)
                if SSE_AVAILABLE:
                    try:
                        sse_manager = get_sse_manager()
                        await sse_manager.broadcast("vectorization_update", {
                            "file_id": file_id,
                            "filename": file_info.filename,
                            "status": "failed",
                            "vectorized": False,
                            "error_message": result.get('error')
                        })
                        self.logger.info(f"📡 SSE 벡터화 실패 이벤트 전송: {file_id}")
                    except Exception as sse_error:
                        self.logger.warning(f"SSE 실패 이벤트 전송 실패: {sse_error}")
                
                return {"success": False, "error": result.get('error')}

        except Exception as e:
            self.logger.error(f"💥 벡터화 오류: {e}", exc_info=True)
            await self._update_file_status(file_id, FileStatus.FAILED, error=str(e))
            
            # SSE 이벤트 전송 (벡터화 예외 오류)
            if SSE_AVAILABLE:
                try:
                    file_info = await self.get_file_info(file_id)
                    sse_manager = get_sse_manager()
                    await sse_manager.broadcast("vectorization_update", {
                        "file_id": file_id,
                        "filename": file_info.filename if file_info else "Unknown",
                        "status": "failed",
                        "vectorized": False,
                        "error_message": str(e)
                    })
                    self.logger.info(f"📡 SSE 벡터화 예외 오류 이벤트 전송: {file_id}")
                except Exception as sse_error:
                    self.logger.warning(f"SSE 예외 오류 이벤트 전송 실패: {sse_error}")
            
            return {"success": False, "error": str(e)}

    # 하위 호환성을 위한 기존 메서드 (deprecated)
    async def start_vectorization_pipeline(self, file_id: str):
        """전체 벡터화 파이프라인을 시작합니다. (deprecated - 전처리와 벡터화가 분리됨)"""
        self.logger.warning(f"start_vectorization_pipeline은 deprecated 됩니다. start_preprocessing과 start_vectorization을 순차적으로 사용하세요.")
        
        # 자동으로 전처리 먼저 실행
        preprocess_result = await self.start_preprocessing(file_id)
        if not preprocess_result.get("success"):
            return preprocess_result
            
        # 전처리 성공 시 벡터화 실행
        return await self.start_vectorization(file_id)

    # 하위 호환성을 위한 preprocess_file 메서드
    async def preprocess_file(self, file_id: str, method: str = "basic"):
        """파일 전처리 - start_preprocessing의 별칭 (하위 호환성)"""
        return await self.start_preprocessing(file_id, method)

    async def _update_file_status(self, file_id: str, status: FileStatus, error: str = None, chunks_count: int = None):
        """파일 상태를 업데이트합니다."""
        # 기존 파일 정보 조회
        file_metadata = self.file_metadata_service.get_file(file_id)
        if not file_metadata:
            self.logger.error(f"파일 메타데이터를 찾을 수 없습니다: {file_id}")
            return
        
        old_status = file_metadata.status
        self.logger.info(f"📝 파일 상태 변경: {old_status} → {status}")
        
        # 업데이트할 필드들
        update_fields = {}
        
        if chunks_count is not None:
            old_chunk_count = file_metadata.chunk_count
            update_fields['chunk_count'] = chunks_count
            self.logger.info(f"📊 청크 수 업데이트: {old_chunk_count} → {chunks_count}개")
        
        if error:
            update_fields['error_message'] = error
            self.logger.error(f"❌ 실패 상태로 변경, 에러 메시지: {error}")
        elif status == FileStatus.COMPLETED:
            update_fields['error_message'] = None
            self.logger.info(f"🔄 에러 메시지 초기화 완료")
        
        # SQLite 상태 업데이트
        success = self.file_metadata_service.update_status(
            file_id=file_id,
            status=status,
            **update_fields
        )
        
        if success:
            self.logger.info(f"✅ SQLite 메타데이터 업데이트 완료")
            
            # 업데이트 후 확인
            updated_file = self.file_metadata_service.get_file(file_id)
            if updated_file and status == FileStatus.COMPLETED:
                self.logger.info(f"🔍 업데이트 후 vectorized 값: {updated_file.vectorized}")
        else:
            self.logger.error(f"❌ SQLite 메타데이터 업데이트 실패: {file_id}")

    async def _save_preprocessed_content(self, file_id: str, text_content: str, method: str):
        """전처리된 텍스트 내용을 파일로 저장합니다."""
        preprocessed_dir = os.path.join(settings.DATA_DIR, "preprocessed")
        os.makedirs(preprocessed_dir, exist_ok=True)
        
        preprocessed_file = os.path.join(preprocessed_dir, f"{file_id}.txt")
        async with aiofiles.open(preprocessed_file, 'w', encoding='utf-8') as f:
            await f.write(text_content)
            
        # SQLite에 전처리 방법 저장
        self.file_metadata_service.update_file(
            file_id=file_id,
            preprocessing_method=method
        )

    async def _load_preprocessed_content(self, file_id: str) -> str:
        """전처리된 텍스트 내용을 로드합니다."""
        preprocessed_file = os.path.join(settings.DATA_DIR, "preprocessed", f"{file_id}.txt")
        if os.path.exists(preprocessed_file):
            async with aiofiles.open(preprocessed_file, 'r', encoding='utf-8') as f:
                return await f.read()
        return ""

    # 하위 호환성을 위한 deprecated 메서드
    async def _update_vectorization_status(self, file_id: str, status: str, error: str = None, chunks_count: int = None):
        """하위 호환성을 위한 deprecated 메서드"""
        self.logger.warning("_update_vectorization_status는 deprecated되었습니다. _update_file_status를 사용하세요.")
        
        # 기존 상태를 새로운 FileStatus로 매핑
        status_mapping = {
            "processing": FileStatus.PREPROCESSING,
            "completed": FileStatus.COMPLETED,
            "failed": FileStatus.FAILED
        }
        new_status = status_mapping.get(status, FileStatus.FAILED)
        await self._update_file_status(file_id, new_status, error, chunks_count)

    async def update_file_vectorization_status(self, file_id: str, vectorized: bool, error_message: str = None, chunk_count: int = None):
        """파일의 벡터화 상태를 업데이트합니다 (API에서 호출용)"""
        try:
            # 업데이트할 필드들
            update_fields = {
                'vectorized': vectorized
            }
            
            if vectorized:
                update_fields['error_message'] = None
                if chunk_count is not None:
                    update_fields['chunk_count'] = chunk_count
                
                # 상태를 COMPLETED로 업데이트
                success = self.file_metadata_service.update_status(
                    file_id=file_id,
                    status=FileStatus.COMPLETED,
                    **update_fields
                )
                
                if success:
                    self.logger.info(f"파일 {file_id} 벡터화 상태를 성공으로 업데이트 (청크: {chunk_count}개)")
                    return True
            else:
                update_fields['error_message'] = error_message
                status = FileStatus.UPLOADED if not error_message else FileStatus.FAILED
                
                success = self.file_metadata_service.update_status(
                    file_id=file_id,
                    status=status,
                    **update_fields
                )
                
                if success:
                    self.logger.info(f"파일 {file_id} 벡터화 상태를 재설정으로 업데이트")
                    return True
            
            self.logger.error(f"파일 메타데이터 업데이트 실패: {file_id}")
            return False
                
        except Exception as e:
            self.logger.error(f"벡터화 상태 업데이트 실패: {file_id}, 오류: {str(e)}")
            return False

    # --- File Management Methods (Preserved) ---
    def _load_allowed_extensions(self):
        try:
            from .settings_service import settings_service
            system_settings = settings_service.get_section_settings("system")
            allowed_file_types = system_settings.get("allowedFileTypes", ["pdf"])
            self.allowed_extensions = {f".{ext}" if not ext.startswith('.') else ext for ext in allowed_file_types}
        except Exception as e:
            self.logger.warning(f"설정 로드 실패, 기본 확장자 사용: {str(e)}")
            self.allowed_extensions = {".pdf", ".docx", ".pptx", ".xlsx", ".txt"}
    
    # FileInfo 속성 하위 호환성 (schema에서 vectorized 속성 요구)
    def _convert_to_file_info(self, file_metadata: FileMetadata) -> FileInfo:
        """FileMetadata를 FileInfo로 변환하는 헬퍼 메서드"""
        return FileInfo(
            file_id=file_metadata.file_id,
            filename=file_metadata.filename,
            saved_filename=file_metadata.saved_filename,
            file_path=file_metadata.file_path,
            file_size=file_metadata.file_size,
            file_hash=file_metadata.file_hash,
            category_id=file_metadata.category_id,
            category_name=file_metadata.category_name,
            status=file_metadata.status,
            upload_time=file_metadata.upload_time,
            preprocessing_started_at=file_metadata.preprocessing_started_at,
            preprocessing_completed_at=file_metadata.preprocessing_completed_at,
            vectorization_started_at=file_metadata.vectorization_started_at,
            vectorization_completed_at=file_metadata.vectorization_completed_at,
            error_message=file_metadata.error_message,
            chunk_count=file_metadata.chunk_count,
            preprocessing_method=file_metadata.preprocessing_method,
            vectorized=file_metadata.vectorized,  # 추가된 필드
            # PDF 변환 정보 필드 추가
            is_converted_to_pdf=file_metadata.is_converted_to_pdf,
            original_extension=file_metadata.original_extension,
            conversion_method=file_metadata.conversion_method
        )

    def _ensure_data_dir(self):
        os.makedirs(settings.DATA_DIR, exist_ok=True)

    async def upload_file(self, file: UploadFile, category_id: Optional[str] = None, allow_global_duplicates: bool = False, force_replace: bool = False, convert_to_pdf: bool = False) -> FileUploadResponse:
        try:
            file_extension = os.path.splitext(file.filename)[1].lower()
            if file_extension not in self.allowed_extensions:
                raise HTTPException(status_code=400, detail=f"지원하지 않는 파일 형식입니다. 지원 형식: {', '.join(self.allowed_extensions)}")
            
            from .settings_service import settings_service
            system_settings = settings_service.get_section_settings("system")
            max_file_size_mb = system_settings.get("maxFileSize", 10)
            max_file_size_bytes = max_file_size_mb * 1024 * 1024
            
            if file.size and file.size > max_file_size_bytes:
                raise HTTPException(status_code=400, detail=f"파일 크기가 너무 큽니다. 최대 크기: {max_file_size_mb}MB")
            
            content = await file.read()
            file_size = len(content)
            
            import hashlib
            file_hash = hashlib.md5(content).hexdigest()
            
            # SQLite에서 중복 파일 검사
            existing_file = self.file_metadata_service.get_file_by_hash(file_hash)
            if existing_file and (allow_global_duplicates or existing_file.category_id == category_id):
                if force_replace:
                    await self.delete_file(existing_file.file_id)
                else:
                    raise HTTPException(
                        status_code=409,
                        detail={"error": "duplicate_file", "message": "동일한 파일이 이미 존재합니다."}
                    )
            
            file_id = str(uuid.uuid4())
            saved_filename = f"{file_id}{file_extension}"
            file_path = os.path.join(self.upload_dir, saved_filename)
            
            category_name = None
            if category_id:
                category = await self.category_service.get_category(category_id)
                if not category:
                    raise HTTPException(status_code=400, detail="존재하지 않는 카테고리입니다.")
                category_name = category.name
            
            async with aiofiles.open(file_path, 'wb') as f:
                await f.write(content)
            
            # PDF 변환 처리 및 변환 정보 저장
            original_extension = None
            is_converted = False
            conversion_method = None
            
            if convert_to_pdf and file_extension.lower() != '.pdf':
                self.logger.info(f"PDF 변환 시작: {file.filename} -> PDF")
                original_extension = file_extension  # 원본 확장자 저장
                try:
                    pdf_path = await self._convert_to_pdf(file_path, file_id)
                    if pdf_path and os.path.exists(pdf_path):
                        # 원본 파일 삭제하고 PDF로 교체
                        os.remove(file_path)
                        file_path = pdf_path
                        saved_filename = f"{file_id}.pdf"
                        file_extension = '.pdf'
                        # 변환된 PDF 파일 크기 재계산
                        file_size = os.path.getsize(file_path)
                        # 변환 성공 정보 저장
                        is_converted = True
                        conversion_method = "python_lib"
                        self.logger.info(f"PDF 변환 완료: {saved_filename}")
                    else:
                        self.logger.error(f"PDF 변환 실패: {file.filename}")
                        raise HTTPException(status_code=500, detail="PDF 변환에 실패했습니다.")
                except Exception as convert_error:
                    self.logger.error(f"PDF 변환 중 오류: {convert_error}")
                    # 변환 실패 시 원본 파일 유지하고 경고만 로그
                    self.logger.warning(f"PDF 변환 실패로 원본 파일 유지: {file.filename}")
                    # 변환 실패 시 변환 정보 초기화
                    original_extension = None
                    is_converted = False
                    conversion_method = None
            
            # 기본 전처리 방법을 설정에서 가져오기
            system_settings = settings_service.get_section_settings("system")
            default_preprocessing_method = system_settings.get("preprocessing_method", "basic")
            
            # SQLite에 파일 메타데이터 저장
            file_metadata = FileMetadata(
                file_id=file_id,
                filename=file.filename,
                saved_filename=saved_filename,
                file_path=file_path,
                file_size=file_size,
                file_hash=file_hash,
                category_id=category_id,
                category_name=category_name,
                status=FileStatus.UPLOADED,
                upload_time=datetime.now(),
                preprocessing_method=default_preprocessing_method,  # 설정에서 가져온 기본값
                # PDF 변환 정보 저장
                is_converted_to_pdf=is_converted,
                original_extension=original_extension,
                conversion_method=conversion_method
            )
            
            success = self.file_metadata_service.create_file(file_metadata)
            if not success:
                # 파일 삭제 후 오류 발생
                if os.path.exists(file_path):
                    os.remove(file_path)
                raise HTTPException(status_code=500, detail="파일 메타데이터 저장에 실패했습니다.")
            
            return FileUploadResponse(
                file_id=file_id,
                filename=file.filename,
                status=FileStatus.UPLOADED,
                file_size=file_size,
                category_id=category_id,
                category_name=category_name,
                message="파일이 성공적으로 업로드되었습니다. 전처리는 별도로 시작해야 합니다."
            )
            
        except HTTPException:
            raise
        except Exception as e:
            self.logger.error(f"파일 업로드 중 오류: {e}", exc_info=True)
            raise HTTPException(status_code=500, detail=f"파일 업로드 중 오류가 발생했습니다: {str(e)}")

    async def list_files(self, category_id: Optional[str] = None) -> List[FileInfo]:
        file_metadatas = self.file_metadata_service.list_files(
            category_id=category_id,
            include_deleted=False
        )
        
        files = []
        for file_metadata in file_metadatas:
            # 파일 존재 여부 확인
            if os.path.exists(file_metadata.file_path):
                files.append(self._convert_to_file_info(file_metadata))
        
        return files

    async def get_file_info(self, file_id: str) -> Optional[FileInfo]:
        file_metadata = self.file_metadata_service.get_file(file_id)
        if not file_metadata:
            return None
        
        return self._convert_to_file_info(file_metadata)
    
    async def get_file_content(self, file_id: str) -> Dict[str, Any]:
        """파일 내용을 추출하여 반환 (청킹용)"""
        try:
            # 파일 정보 조회
            file_metadata = self.file_metadata_service.get_file(file_id)
            if not file_metadata:
                return {"success": False, "error": "파일을 찾을 수 없습니다"}
            
            # 전처리된 텍스트 파일 경로
            preprocessed_path = os.path.join(settings.DATA_DIR, "preprocessed", f"{file_id}.txt")
            
            # 전처리된 파일이 있으면 해당 내용 사용
            if os.path.exists(preprocessed_path):
                try:
                    with open(preprocessed_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                    return {"success": True, "content": content}
                except Exception as e:
                    self.logger.error(f"전처리된 파일 읽기 실패 {file_id}: {e}")
            
            # 전처리된 파일이 없으면 원본 파일에서 추출
            original_path = file_metadata.file_path
            if not os.path.exists(original_path):
                return {"success": False, "error": "원본 파일을 찾을 수 없습니다"}
            
            # 실제 파일 경로의 확장자로 처리 방식 결정 (변환된 파일 지원)
            file_path_lower = original_path.lower()
            
            # 텍스트 파일인 경우 직접 읽기
            if file_path_lower.endswith(('.txt', '.md', '.html', '.json', '.xml', '.csv')):
                try:
                    with open(original_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                    return {"success": True, "content": content}
                except UnicodeDecodeError:
                    try:
                        with open(original_path, 'r', encoding='cp949') as f:
                            content = f.read()
                        return {"success": True, "content": content}
                    except Exception as e:
                        return {"success": False, "error": f"텍스트 파일 읽기 실패: {str(e)}"}
            
            # PPTX 파일인 경우 - 텍스트 추출 (빠른청킹용)
            elif file_path_lower.endswith('.pptx'):
                try:
                    from pptx import Presentation
                    
                    prs = Presentation(original_path)
                    text_content = ""
                    
                    for i, slide in enumerate(prs.slides, 1):
                        slide_text = f"\n=== 슬라이드 {i} ===\n"
                        
                        # 슬라이드의 모든 텍스트 추출
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text += shape.text + "\n"
                        
                        text_content += slide_text
                    
                    return {"success": True, "content": text_content.strip()}
                    
                except ImportError:
                    return {"success": False, "error": "PPTX 텍스트 추출을 위해 python-pptx 패키지가 필요합니다."}
                except Exception as e:
                    self.logger.error(f"PPTX 텍스트 추출 실패 {file_id}: {e}")
                    return {"success": False, "error": f"PPTX 텍스트 추출 중 오류 발생: {str(e)}"}
            
            # DOCX 파일인 경우 - 텍스트 추출 (빠른청킹용)
            elif file_path_lower.endswith('.docx'):
                try:
                    from docx import Document
                    
                    doc = Document(original_path)
                    text_content = ""
                    
                    for paragraph in doc.paragraphs:
                        if paragraph.text.strip():
                            text_content += paragraph.text + "\n"
                    
                    return {"success": True, "content": text_content.strip()}
                    
                except ImportError:
                    return {"success": False, "error": "DOCX 텍스트 추출을 위해 python-docx 패키지가 필요합니다."}
                except Exception as e:
                    self.logger.error(f"DOCX 텍스트 추출 실패 {file_id}: {e}")
                    return {"success": False, "error": f"DOCX 텍스트 추출 중 오류 발생: {str(e)}"}
            
            # PDF 파일인 경우 - 텍스트 추출 (빠른청킹용)
            elif file_path_lower.endswith('.pdf'):
                try:
                    import fitz  # PyMuPDF
                    
                    doc = fitz.open(original_path)
                    text_content = ""
                    
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        page_text = page.get_text()
                        if page_text.strip():
                            text_content += f"\n=== 페이지 {page_num + 1} ===\n"
                            text_content += page_text + "\n"
                    
                    doc.close()
                    return {"success": True, "content": text_content.strip()}
                    
                except ImportError:
                    return {"success": False, "error": "PDF 텍스트 추출을 위해 PyMuPDF 패키지가 필요합니다."}
                except Exception as e:
                    self.logger.error(f"PDF 텍스트 추출 실패 {file_id}: {e}")
                    return {"success": False, "error": f"PDF 텍스트 추출 중 오류 발생: {str(e)}"}
            
            # 기타 파일은 전처리가 필요
            else:
                return {"success": False, "error": "파일이 아직 전처리되지 않았습니다. 먼저 전처리를 진행해주세요."}
            
        except Exception as e:
            self.logger.error(f"파일 내용 추출 실패 {file_id}: {e}")
            return {"success": False, "error": f"파일 내용 추출 중 오류 발생: {str(e)}"}

    async def delete_file(self, file_id: str) -> bool:
        file_metadata = self.file_metadata_service.get_file(file_id)
        if not file_metadata:
            return False
        
        try:
            # 물리 파일 삭제
            if os.path.exists(file_metadata.file_path):
                os.remove(file_metadata.file_path)
            
            # 벡터 데이터 삭제
            await self.vector_service.delete_document_vectors(file_id)
            
            # SQLite에서 소프트 삭제 (status를 deleted로 변경)
            success = self.file_metadata_service.delete_file(file_id, soft_delete=True)
            return success
        except Exception as e:
            self.logger.error(f"파일 삭제 실패: {e}", exc_info=True)
            return False
    
    async def get_file_path(self, file_id: str) -> Optional[str]:
        """파일 경로 조회"""
        file_metadata = self.file_metadata_service.get_file(file_id)
        if file_metadata and os.path.exists(file_metadata.file_path):
            return file_metadata.file_path
        return None
    
    
    
    async def retry_vectorization(self, file_id: str) -> bool:
        """벡터화 재시도"""
        try:
            file_metadata = self.file_metadata_service.get_file(file_id)
            if not file_metadata:
                return False
            
            # 상태를 UPLOADED로 재설정하고 벡터화 재시도
            self.file_metadata_service.update_status(
                file_id=file_id,
                status=FileStatus.UPLOADED,
                vectorized=False,
                error_message=None
            )
            
            # 백그라운드에서 벡터화 시작
            import asyncio
            asyncio.create_task(self.start_vectorization(file_id))
            
            return True
            
        except Exception as e:
            self.logger.error(f"벡터화 재시도 실패: {e}")
            return False
    
    async def set_search_flow(self, flow_id: str) -> bool:
        """검색 Flow 설정 (하위 호환성)"""
        self.logger.warning("set_search_flow는 현재 구현되지 않았습니다.")
        return False
    
    async def delete_flow(self, flow_id: str) -> bool:
        """Flow 삭제 (하위 호환성)"""
        self.logger.warning("delete_flow는 현재 구현되지 않았습니다.")
        return False
    
    async def cleanup_orphaned_metadata(self) -> int:
        """고아 메타데이터 정리"""
        try:
            files = self.file_metadata_service.list_files(include_deleted=False)
            orphaned_count = 0
            
            for file_metadata in files:
                # 실제 파일이 존재하지 않는 메타데이터를 찾음
                if not os.path.exists(file_metadata.file_path):
                    success = self.file_metadata_service.delete_file(
                        file_id=file_metadata.file_id,
                        soft_delete=True
                    )
                    if success:
                        orphaned_count += 1
                        self.logger.info(f"고아 메타데이터 삭제: {file_metadata.file_id}")
            
            return orphaned_count
            
        except Exception as e:
            self.logger.error(f"고아 메타데이터 정리 실패: {e}")
            return 0
    
    async def sync_vectorization_status(self) -> Dict[str, Any]:
        """벡터화 상태 동기화"""
        try:
            files = self.file_metadata_service.list_files(include_deleted=False)
            corrected_count = 0
            
            for file_metadata in files:
                # 벡터 서비스에서 실제 벡터 데이터 존재 여부 확인
                try:
                    has_vectors = await self.vector_service.has_document_vectors(file_metadata.file_id)
                    
                    # 메타데이터와 실제 상태가 다른 경우 수정
                    if file_metadata.vectorized != has_vectors:
                        status = FileStatus.COMPLETED if has_vectors else FileStatus.UPLOADED
                        self.file_metadata_service.update_status(
                            file_id=file_metadata.file_id,
                            status=status,
                            vectorized=has_vectors
                        )
                        corrected_count += 1
                        self.logger.info(f"벡터화 상태 동기화: {file_metadata.file_id} -> {has_vectors}")
                        
                except Exception as e:
                    self.logger.warning(f"벡터 상태 확인 실패 {file_metadata.file_id}: {e}")
            
            return {
                "status_corrected": corrected_count,
                "message": f"{corrected_count}개 파일의 상태를 동기화했습니다."
            }
            
        except Exception as e:
            self.logger.error(f"벡터화 상태 동기화 실패: {e}")
            return {"status_corrected": 0, "error": str(e)}
    
    async def _register_korean_font(self) -> str:
        """한글 폰트를 등록하고 폰트명을 반환합니다."""
        try:
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            # Windows 시스템 폰트 경로에서 한글 폰트 찾기
            font_paths = [
                "C:/Windows/Fonts/malgun.ttf",  # 맑은고딕
                "C:/Windows/Fonts/batang.ttf",   # 바탕
                "C:/Windows/Fonts/gulim.ttf",    # 굴림
            ]
            
            for font_path in font_paths:
                if os.path.exists(font_path):
                    font_name = os.path.basename(font_path).split('.')[0]
                    try:
                        pdfmetrics.registerFont(TTFont(font_name, font_path))
                        self.logger.info(f"한글 폰트 등록 성공: {font_name}")
                        return font_name
                    except Exception as font_error:
                        self.logger.warning(f"폰트 등록 실패 {font_name}: {font_error}")
                        continue
            
            self.logger.warning("한글 폰트를 찾을 수 없습니다. 기본 폰트를 사용합니다.")
            return "Helvetica"
            
        except Exception as e:
            self.logger.error(f"폰트 등록 중 오류: {e}")
            return "Helvetica"
    
    async def _convert_to_pdf(self, file_path: str, file_id: str) -> Optional[str]:
        """파일을 PDF로 변환합니다."""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            import openpyxl
            from docx import Document
            
            # 변환된 PDF 파일 경로
            pdf_path = os.path.join(self.upload_dir, f"{file_id}.pdf")
            file_extension = os.path.splitext(file_path)[1].lower()
            
            self.logger.info(f"Python 라이브러리를 사용하여 PDF 변환 시작: {file_path} ({file_extension})")
            
            # 파일 형식별 변환 로직
            if file_extension == '.xlsx' or file_extension == '.xls':
                # Excel 파일 변환
                return await self._convert_excel_to_pdf(file_path, pdf_path)
            elif file_extension == '.docx' or file_extension == '.doc':
                # Word 파일 변환
                return await self._convert_word_to_pdf(file_path, pdf_path)
            elif file_extension == '.txt':
                # 텍스트 파일 변환
                return await self._convert_text_to_pdf(file_path, pdf_path)
            else:
                self.logger.warning(f"지원되지 않는 변환 형식: {file_extension}")
                return None
                
        except ImportError as e:
            self.logger.error(f"필요한 라이브러리가 설치되지 않았습니다: {e}")
            self.logger.info("pip install reportlab openpyxl python-docx 명령으로 필요한 라이브러리를 설치해주세요.")
            return None
        except Exception as e:
            self.logger.error(f"PDF 변환 중 예외 발생: {e}")
            return None
    
    async def _convert_excel_to_pdf(self, excel_path: str, pdf_path: str) -> Optional[str]:
        """Excel 파일을 PDF로 변환"""
        try:
            import openpyxl
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            # 한글 폰트 등록
            korean_font = await self._register_korean_font()
            
            # Excel 파일 읽기 (data_only=True로 수식의 계산된 값 가져오기)
            workbook = openpyxl.load_workbook(excel_path, data_only=True)
            sheet = workbook.active
            
            # PDF 생성
            c = canvas.Canvas(pdf_path, pagesize=letter)
            width, height = letter
            
            y_position = height - 50
            
            # 제목 추가
            c.setFont(korean_font, 16)
            title = f"Excel 파일 변환: {os.path.basename(excel_path)}"
            c.drawString(50, y_position, title)
            y_position -= 40
            
            # 셀 데이터 출력 (셀 객체에 직접 접근)
            c.setFont(korean_font, 10)
            for row in sheet.iter_rows():
                if y_position < 50:  # 새 페이지
                    c.showPage()
                    c.setFont(korean_font, 10)  # 새 페이지에서도 폰트 재설정
                    y_position = height - 50
                
                # 한글 텍스트 처리
                row_cells = []
                for cell in row:
                    if cell.value is not None:
                        # 계산된 값 우선 사용
                        if hasattr(cell, 'displayed_value') and cell.displayed_value is not None:
                            cell_str = str(cell.displayed_value)
                        elif cell.data_type == 'f' and cell.value:  # 수식 셀
                            # 수식의 계산된 결과값 시도
                            try:
                                # 계산된 값이 있으면 사용
                                if hasattr(cell, '_value') and cell._value is not None:
                                    cell_str = str(cell._value)
                                else:
                                    cell_str = f"[계산값 없음: {str(cell.value)[:15]}...]"
                            except:
                                cell_str = f"[수식: {str(cell.value)[:15]}...]"
                        else:
                            cell_str = str(cell.value)
                        
                        # 한글이 포함된 텍스트는 더 짧게 자르기
                        if any('\u4e00' <= char <= '\u9fff' or '\uac00' <= char <= '\ud7af' for char in cell_str):
                            cell_str = cell_str[:30]  # 한글/한자 포함 시 짧게
                        else:
                            cell_str = cell_str[:50]  # 영문은 좀 더 길게
                        row_cells.append(cell_str)
                    else:
                        row_cells.append("")
                
                row_text = " | ".join(row_cells)
                try:
                    c.drawString(50, y_position, row_text)
                except Exception as draw_error:
                    # 한글 출력 실패 시 ASCII만 출력
                    ascii_text = ''.join(char if ord(char) < 128 else '?' for char in row_text)
                    c.drawString(50, y_position, ascii_text)
                    
                y_position -= 15
            
            c.save()
            self.logger.info(f"Excel -> PDF 변환 완료: {pdf_path}")
            return pdf_path
            
        except Exception as e:
            self.logger.error(f"Excel PDF 변환 실패: {e}")
            return None
    
    async def _convert_word_to_pdf(self, word_path: str, pdf_path: str) -> Optional[str]:
        """Word 파일을 PDF로 변환"""
        try:
            from docx import Document
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            # 한글 폰트 등록
            korean_font = await self._register_korean_font()
            
            # Word 파일 읽기
            doc = Document(word_path)
            
            # PDF 생성
            c = canvas.Canvas(pdf_path, pagesize=letter)
            width, height = letter
            
            y_position = height - 50
            
            # 제목 추가
            c.setFont(korean_font, 16)
            title = f"Word 파일 변환: {os.path.basename(word_path)}"
            c.drawString(50, y_position, title)
            y_position -= 40
            
            # 문단 출력
            c.setFont(korean_font, 12)
            for paragraph in doc.paragraphs:
                if y_position < 50:  # 새 페이지
                    c.showPage()
                    c.setFont(korean_font, 12)
                    y_position = height - 50
                
                # 긴 텍스트를 여러 줄로 나누기
                text = paragraph.text
                if text.strip():
                    # 한글 고려하여 줄 길이 조정
                    line_length = 40 if any('\uac00' <= char <= '\ud7af' for char in text) else 80
                    lines = [text[i:i+line_length] for i in range(0, len(text), line_length)]
                    for line in lines:
                        if y_position < 50:
                            c.showPage()
                            c.setFont(korean_font, 12)
                            y_position = height - 50
                        try:
                            c.drawString(50, y_position, line)
                        except Exception:
                            # 한글 출력 실패 시 ASCII만 출력
                            ascii_line = ''.join(char if ord(char) < 128 else '?' for char in line)
                            c.drawString(50, y_position, ascii_line)
                        y_position -= 15
                    y_position -= 5  # 문단 간격
            
            c.save()
            self.logger.info(f"Word -> PDF 변환 완료: {pdf_path}")
            return pdf_path
            
        except Exception as e:
            self.logger.error(f"Word PDF 변환 실패: {e}")
            return None
    
    async def _convert_text_to_pdf(self, text_path: str, pdf_path: str) -> Optional[str]:
        """텍스트 파일을 PDF로 변환"""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            # 한글 폰트 등록
            korean_font = await self._register_korean_font()
            
            # 텍스트 파일 읽기
            with open(text_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
            
            # PDF 생성
            c = canvas.Canvas(pdf_path, pagesize=letter)
            width, height = letter
            
            y_position = height - 50
            
            # 제목 추가
            c.setFont(korean_font, 16)
            title = f"텍스트 파일 변환: {os.path.basename(text_path)}"
            c.drawString(50, y_position, title)
            y_position -= 40
            
            # 텍스트 출력
            c.setFont(korean_font, 12)
            lines = text_content.split('\n')
            for line in lines:
                if y_position < 50:  # 새 페이지
                    c.showPage()
                    c.setFont(korean_font, 12)
                    y_position = height - 50
                
                # 한글 고려하여 줄 길이 조정
                line_length = 40 if any('\uac00' <= char <= '\ud7af' for char in line) else 80
                
                # 긴 줄을 여러 줄로 나누기
                if len(line) > line_length:
                    sub_lines = [line[i:i+line_length] for i in range(0, len(line), line_length)]
                    for sub_line in sub_lines:
                        if y_position < 50:
                            c.showPage()
                            c.setFont(korean_font, 12)
                            y_position = height - 50
                        try:
                            c.drawString(50, y_position, sub_line)
                        except Exception:
                            # 한글 출력 실패 시 ASCII만 출력
                            ascii_line = ''.join(char if ord(char) < 128 else '?' for char in sub_line)
                            c.drawString(50, y_position, ascii_line)
                        y_position -= 15
                else:
                    try:
                        c.drawString(50, y_position, line)
                    except Exception:
                        # 한글 출력 실패 시 ASCII만 출력
                        ascii_line = ''.join(char if ord(char) < 128 else '?' for char in line)
                        c.drawString(50, y_position, ascii_line)
                    y_position -= 15
            
            c.save()
            self.logger.info(f"Text -> PDF 변환 완료: {pdf_path}")
            return pdf_path
            
        except Exception as e:
            self.logger.error(f"Text PDF 변환 실패: {e}")
            return None
