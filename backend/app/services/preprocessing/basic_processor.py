
import os
import logging
import aiofiles
import json
import re
from typing import Dict, Any

# 로거 설정
logger = logging.getLogger(__name__)

# PDF 처리
async def _extract_pdf_with_pymupdf(file_path: str) -> str:
    """pymupdf를 사용하여 PDF 텍스트 추출 (CID 깨짐 방지 강화)"""
    import fitz  # pymupdf
    logger.info("📄 (Basic) pymupdf로 PDF 텍스트 추출 시도 (CID 처리 강화)...")
    doc = fitz.open(file_path)
    extracted_text = ""
    for page_num in range(len(doc)):
        try:
            page = doc.load_page(page_num)
            page_text = page.get_text("text")
            cleaned_text = re.sub(r'\(cid:\d+\)', ' ', page_text)
            if cleaned_text.strip():
                extracted_text += f"[페이지 {page_num + 1}]\n{cleaned_text}\n\n"
        except Exception as e:
            logger.warning(f"페이지 {page_num + 1} 처리 실패: {str(e)}")
    doc.close()
    if extracted_text.strip():
        logger.info("✅ (Basic) pymupdf로 PDF 텍스트 추출 성공")
        return extracted_text.strip()
    else:
        raise Exception("pymupdf 텍스트 추출 결과가 비어있음")

async def _extract_pdf_with_pypdf(file_path: str) -> str:
    """pypdf를 사용하여 PDF 텍스트 추출"""
    import pypdf
    logger.info("📄 (Basic) pypdf로 PDF 텍스트 추출 시도...")
    extracted_text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = pypdf.PdfReader(file)
        for page_num, page in enumerate(pdf_reader.pages):
            try:
                page_text = page.extract_text()
                if page_text.strip():
                    extracted_text += f"[페이지 {page_num + 1}]\n{page_text}\n\n"
            except Exception as e:
                logger.warning(f"페이지 {page_num + 1} 처리 실패: {str(e)}")
    if extracted_text.strip():
        logger.info("✅ (Basic) pypdf로 PDF 텍스트 추출 성공")
        return extracted_text.strip()
    else:
        raise Exception("pypdf 텍스트 추출 결과가 비어있음")

async def _extract_pdf_with_pdfminer(file_path: str) -> str:
    """pdfminer를 사용하여 PDF 텍스트 추출"""
    from pdfminer.high_level import extract_text
    logger.info("📄 (Basic) pdfminer로 PDF 텍스트 추출 시도...")
    extracted_text = extract_text(file_path)
    if extracted_text.strip():
        logger.info("✅ (Basic) pdfminer로 PDF 텍스트 추출 성공")
        return extracted_text.strip()
    else:
        raise Exception("pdfminer 텍스트 추출 결과가 비어있음")

async def _extract_pdf_fallback(file_path: str) -> str:
    """PDF 파일에 대한 폴백 처리"""
    # 설정은 중앙 설정 서비스에서 받아오도록 수정될 예정
    fallback_order = ["pypdf", "pymupdf", "pdfminer"]
    logger.info(f"PDF 폴백 처리 시작 (순서: {fallback_order})")
    for processor in fallback_order:
        try:
            if processor == "pymupdf":
                return await _extract_pdf_with_pymupdf(file_path)
            elif processor == "pypdf":
                return await _extract_pdf_with_pypdf(file_path)
            elif processor == "pdfminer":
                return await _extract_pdf_with_pdfminer(file_path)
        except Exception as e:
            logger.warning(f"{processor} 처리 실패: {str(e)}")
            continue
    raise Exception("모든 PDF 처리 방법이 실패했습니다")

# Office 파일 처리
async def _extract_office_text(file_path: str, file_extension: str) -> str:
    text = ""
    if file_extension in ['.doc', '.docx']:
        try:
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
        except Exception as e:
            return f"⚠️ Word 파일 처리 실패: {e}"
    elif file_extension in ['.ppt', '.pptx']:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text += shape.text + "\n"
        except Exception as e:
            return f"⚠️ PowerPoint 파일 처리 실패: {e}"
    elif file_extension in ['.xls', '.xlsx']:
        try:
            import pandas as pd
            # Excel 파일의 모든 시트를 읽어서 텍스트로 변환
            excel_file = pd.ExcelFile(file_path)
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    # 시트 이름 추가
                    text += f"\n[시트: {sheet_name}]\n"
                    # 데이터프레임을 텍스트로 변환 (NaN 값 제거)
                    df_text = df.fillna('').to_string(index=False)
                    text += df_text + "\n\n"
                except Exception as sheet_error:
                    logger.warning(f"시트 '{sheet_name}' 처리 실패: {sheet_error}")
                    text += f"\n[시트: {sheet_name}] - 처리 실패: {sheet_error}\n"
        except Exception as e:
            return f"⚠️ Excel 파일 처리 실패: {e}"
    # ... (다른 office 형식 추가) ...
    return text.strip()

# 기타 파일 형식 처리
async def _extract_text_from_text_file(file_path: str) -> str:
    encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
    for enc in encodings:
        try:
            async with aiofiles.open(file_path, 'r', encoding=enc) as f:
                return await f.read()
        except UnicodeDecodeError:
            continue
    raise Exception(f"{file_path}의 인코딩을 확인할 수 없습니다.")

# ... (html, json, xml 등 다른 추출기들) ...

# --- Main Entry Point --- #

async def process(file_path: str, file_extension: str) -> str:
    """파일 확장자에 따라 적절한 기본 텍스트 추출기를 호출합니다."""
    ext_lower = file_extension.lower()
    logger.info(f"⚙️ Basic Processor 실행: {ext_lower}")
    
    if ext_lower == '.pdf':
        return await _extract_pdf_fallback(file_path)
    
    elif ext_lower in ['.doc', '.docx', '.ppt', '.pptx', '.xls', '.xlsx']:
        return await _extract_office_text(file_path, ext_lower)

    elif ext_lower in ['.txt', '.md', '.csv']:
        return await _extract_text_from_text_file(file_path)
    
    # elif ext_lower in ['.html', '.htm']:
    #     return await _extract_text_from_html(file_path)
    
    # ... (다른 형식들) ...

    else:
        logger.warning(f"지원하지 않는 파일 형식 ({ext_lower})에 대한 기본 처리기가 없습니다.")
        raise NotImplementedError(f"Basic processor not implemented for {ext_lower}")
