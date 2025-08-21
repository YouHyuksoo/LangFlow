from fastapi import APIRouter, HTTPException, UploadFile, File, Form, Query, Body
from typing import List, Optional
from ..core.config import settings
from ..core.logger import get_user_logger, get_console_logger
from ..models.schemas import FileUploadResponse, FileInfo, DoclingOptions
from ..services import get_file_service
from ..services.settings_service import settings_service
import os
import json

router = APIRouter(prefix="/files", tags=["files"])
_ulog = get_user_logger()
_clog = get_console_logger()

def get_file_service_instance():
    """FileService 인스턴스를 지연 로딩으로 가져옵니다."""
    return get_file_service()


@router.post("/upload", response_model=FileUploadResponse)
async def upload_file(
    file: UploadFile = File(...),
    category_id: Optional[str] = Form(None),
    category: Optional[str] = Form(None),
    force_replace: bool = Form(False),
    convert_to_pdf: Optional[str] = Form(None)
):
    """파일 업로드"""
    try:
        # convert_to_pdf 문자열을 boolean으로 변환
        convert_pdf_bool = convert_to_pdf == "true" if convert_to_pdf else False
        
        _clog.info(f"파일 업로드 요청 수신: filename={file.filename}, category_id={category_id}, category={category}, force_replace={force_replace}, convert_to_pdf={convert_pdf_bool}")
        
        # 카테고리 검증
        if not category_id and not category:
            _clog.error("카테고리 누락된 업로드 요청")
            raise HTTPException(
                status_code=400, 
                detail="카테고리를 지정해야 합니다. category_id 또는 category를 제공해주세요."
            )
        
        # 파일 형식 검증 - 통합 설정 사용
        system_settings = settings_service.get_section_settings("system")
        allowed_file_types = system_settings.get("allowedFileTypes", ["pdf"])
        allowed_extensions = [f".{ext}" if not ext.startswith('.') else ext for ext in allowed_file_types]
        
        file_extension = None
        if file.filename:
            file_extension = '.' + file.filename.lower().split('.')[-1]
        
        if not file_extension or file_extension not in allowed_extensions:
            _clog.error(f"잘못된 파일 형식: {file.filename}")
            raise HTTPException(
                status_code=400,
                detail=f"지원되지 않는 파일 형식입니다. 지원 형식: {', '.join(allowed_extensions)}"
            )
        
        # 파일 크기 로깅 (검증은 FileService에서 처리)
        _clog.info(f"파일 크기: {file.size} bytes")
        
        
        # FileService 업로드 시도 (자동 전처리 없이 업로드만)
        _clog.info("FileService.upload_file 호출 시작")
        response = await get_file_service_instance().upload_file(
            file, 
            category_id or category, 
            force_replace=force_replace,
            convert_to_pdf=convert_pdf_bool
        )
        _clog.info("FileService.upload_file 호출 완료")
        
        _ulog.info(
            "파일 업로드 완료",
            extra={
                "event": "file_uploaded",
                "file_id": response.file_id,
                "category": category_id or category,
                "uploaded_filename": response.filename,
            },
        )
        return response
    except HTTPException as he:
        _clog.error(f"HTTP 예외 발생: {he.status_code} - {he.detail}")
        raise
    except Exception as e:
        _clog.exception(f"파일 업로드 중 예기치 못한 오류: {str(e)}")
        raise HTTPException(status_code=500, detail=f"파일 업로드 중 오류가 발생했습니다: {str(e)}")

@router.get("/", response_model=List[FileInfo])
async def list_files(category_id: Optional[str] = Query(None)):
    """업로드된 파일 목록 조회"""
    try:
        _clog.debug(f"파일 목록 조회 - category_id: {category_id}")
        files = await get_file_service_instance().list_files(category_id)
        _clog.debug(f"파일 목록 조회 완료 - {len(files)}개")
        return files
    except Exception as e:
        _clog.exception("파일 목록 조회 중 오류", extra={"event": "file_list_error"})
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/category/{category_id}", response_model=List[FileInfo])
async def get_files_by_category(category_id: str):
    """카테고리별 파일 목록 조회"""
    try:
        files = await get_file_service_instance().list_files(category_id)
        return files
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/{file_id}", response_model=FileInfo)
async def get_file_info(file_id: str):
    """특정 파일 정보 조회"""
    try:
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        return file_info
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/{file_id}")
async def delete_file(file_id: str):
    """파일 삭제"""
    try:
        success = await get_file_service_instance().delete_file(file_id)
        if not success:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        _ulog.info("파일 삭제", extra={"event": "file_deleted", "file_id": file_id})
        return {"message": "파일이 성공적으로 삭제되었습니다."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/{file_id}/preprocess")
async def preprocess_file(file_id: str, method: str = "basic"):
    """파일 전처리 (백그라운드 처리)"""
    try:
        # 파일 존재 여부 확인
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        # 백그라운드에서 전처리 시작 (즉시 응답 반환)
        import asyncio
        asyncio.create_task(get_file_service_instance().start_preprocessing(file_id, method))
        
        _ulog.info("파일 전처리 시작", extra={"event": "file_preprocessing_started", "file_id": file_id, "method": method})
        return {"message": f"'{file_info.filename}' 파일의 전처리가 시작되었습니다. (방법: {method})"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# 이 API는 vectorize-with-docling으로 대체되었습니다.

@router.post("/{file_id}/process-complete")
async def process_complete_pipeline(file_id: str, method: str = "basic"):
    """전처리부터 벡터화까지 전체 파이프라인 실행 (백그라운드 처리)"""
    try:
        # 파일 존재 여부 확인
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        # 백그라운드에서 전체 파이프라인 시작 (즉시 응답 반환)
        import asyncio
        asyncio.create_task(get_file_service_instance().start_vectorization(file_id))
        
        _ulog.info("파일 전체 처리 시작", extra={"event": "file_complete_processing_started", "file_id": file_id})
        return {"message": f"'{file_info.filename}' 파일의 전체 처리가 시작되었습니다. (전처리 방법: {method})"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# 이 API도 vectorize-with-docling으로 대체되었습니다.

@router.post("/{file_id}/force-reprocess")
async def force_reprocess_file(file_id: str):
    """PREPROCESSING나 FAILED 상태의 파일을 강제로 재처리합니다."""
    try:
        _clog.info(f"파일 강제 재처리 요청: file_id={file_id}")
        
        # 파일 정보 확인
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        _clog.info(f"현재 파일 상태: {file_info.status}")
        
        # 강제 재처리 가능한 상태인지 확인
        from ..models.schemas import FileStatus
        if file_info.status not in [FileStatus.PREPROCESSING, FileStatus.FAILED, FileStatus.UPLOADED]:
            raise HTTPException(
                status_code=400, 
                detail=f"현재 상태({file_info.status})에서는 강제 재처리를 할 수 없습니다. PREPROCESSING, FAILED, UPLOADED 상태에서만 가능합니다."
            )
        
        file_service = get_file_service_instance()
        
        # 기존 벡터 데이터가 있다면 삭제
        if file_service.vector_service:
            try:
                await file_service.vector_service.delete_document_vectors(file_id)
                _clog.info(f"기존 벡터 데이터 삭제 완료: {file_id}")
            except Exception as e:
                _clog.warning(f"기존 벡터 데이터 삭제 중 오류 (계속 진행): {str(e)}")
        
        # 파일 상태를 UPLOADED로 재설정
        await file_service._update_file_status(file_id, FileStatus.UPLOADED)
        _clog.info(f"파일 상태를 UPLOADED로 재설정: {file_id}")
        
        # 오류 메시지 초기화
        await file_service.update_file_vectorization_status(
            file_id=file_id,
            vectorized=False,
            error_message=None,
            chunk_count=None
        )
        
        # 백그라운드에서 벡터화 시작
        import asyncio
        asyncio.create_task(file_service.start_vectorization(file_id))
        
        _ulog.info("파일 강제 재처리 시작", extra={"event": "file_force_reprocess_started", "file_id": file_id})
        return {
            "message": f"'{file_info.filename}' 파일의 강제 재처리가 시작되었습니다. 기존 데이터가 삭제되고 처음부터 다시 처리됩니다.",
            "previous_status": str(file_info.status),
            "new_status": "UPLOADED"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        _clog.error(f"파일 강제 재처리 중 오류: {str(e)}")
        raise HTTPException(status_code=500, detail=f"강제 재처리 중 오류가 발생했습니다: {str(e)}")

@router.post("/{file_id}/vectorize")
async def vectorize_file(
    file_id: str,
    enable_docling: bool = Body(True),
    extract_tables: bool = Body(True),
    extract_images: bool = Body(True),
    ocr_enabled: bool = Body(False),
    output_format: str = Body("markdown")
):
    """파일 벡터화 (이미지 메타데이터 포함)"""
    try:
        _clog.info(f"Docling 통합 벡터화 요청: file_id={file_id}, enable_docling={enable_docling}")
        
        # 파일 정보 확인
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        # VectorService 인스턴스 가져오기
        file_service = get_file_service_instance()
        if not file_service.vector_service:
            raise HTTPException(status_code=500, detail="VectorService를 사용할 수 없습니다.")
        
        # 기존 벡터 데이터 삭제 (있는 경우)
        try:
            # delete_document_vectors 메서드 존재 여부 확인
            if hasattr(file_service.vector_service, 'delete_document_vectors'):
                await file_service.vector_service.delete_document_vectors(file_id)
                _clog.info(f"기존 벡터 데이터 삭제 완료: {file_id}")
            else:
                # 메서드가 없는 경우 대체 방법 사용
                _clog.warning(f"delete_document_vectors 메서드가 없음. 대체 방법 사용")
                # ChromaDB에 직접 접근하여 삭제 시도
                if hasattr(file_service.vector_service, '_collection') and file_service.vector_service._collection:
                    try:
                        existing_data = file_service.vector_service._collection.get(
                            where={"file_id": file_id},
                            include=["metadatas"]
                        )
                        if existing_data and existing_data['ids']:
                            file_service.vector_service._collection.delete(ids=existing_data['ids'])
                            _clog.info(f"대체 방법으로 벡터 데이터 삭제 완료: {file_id}")
                    except Exception as inner_e:
                        _clog.warning(f"대체 방법으로도 삭제 실패: {str(inner_e)}")
        except Exception as e:
            _clog.warning(f"기존 벡터 데이터 삭제 중 오류 (계속 진행): {str(e)}")
        
        # Docling 옵션 설정
        docling_options = None
        if enable_docling:
            docling_options = DoclingOptions(
                output_format=output_format,
                extract_tables=extract_tables,
                extract_images=extract_images,
                ocr_enabled=ocr_enabled
            )
        
        # 벡터화 메타데이터 준비
        vector_metadata = {
            "filename": file_info.filename,
            "category_id": file_info.category_id,
            "category_name": file_info.category_name,
            "upload_time": file_info.upload_time.isoformat() if file_info.upload_time else None,
            "file_size": file_info.file_size
        }
        
        _clog.info(f"Docling 통합 벡터화 시작: {file_info.filename}")
        
        # 통합 벡터화 파이프라인 실행
        result = await file_service.vector_service.vectorize_with_docling_pipeline(
            file_path=file_info.file_path,
            file_id=file_id,
            metadata=vector_metadata,
            enable_docling=enable_docling,
            docling_options=docling_options
        )
        
        if result["success"]:
            # 파일 상태 업데이트
            await file_service.update_file_vectorization_status(
                file_id=file_id,
                vectorized=True,
                error_message=None,
                chunk_count=result["chunks_count"]
            )
            
            _ulog.info(
                "Docling 통합 벡터화 완료",
                extra={
                    "event": "docling_vectorization_completed",
                    "file_id": file_id,
                    "chunks_count": result["chunks_count"],
                    "processing_method": result.get("processing_method"),
                    "processing_time": result.get("processing_time")
                }
            )
            
            return {
                "success": True,
                "message": f"'{file_info.filename}' 파일이 성공적으로 벡터화되었습니다.",
                "chunks_count": result["chunks_count"],
                "processing_method": result.get("processing_method", "unknown"),
                "processing_time": result.get("processing_time", 0),
                "docling_used": enable_docling
            }
        else:
            # 벡터화 실패 상태 업데이트
            await file_service.update_file_vectorization_status(
                file_id=file_id,
                vectorized=False,
                error_message=result.get("error", "벡터화 실패"),
                chunk_count=0
            )
            
            _clog.error(f"Docling 통합 벡터화 실패: {result.get('error')}")
            
            return {
                "success": False,
                "message": f"'{file_info.filename}' 파일 벡터화에 실패했습니다.",
                "error": result.get("error", "알 수 없는 오류"),
                "docling_used": enable_docling
            }
            
    except HTTPException:
        raise
    except Exception as e:
        _clog.exception(f"Docling 통합 벡터화 중 예기치 못한 오류: {str(e)}")
        raise HTTPException(status_code=500, detail=f"벡터화 중 오류가 발생했습니다: {str(e)}")

@router.post("/{file_id}/revectorize")
async def revectorize_file(
    file_id: str,
    enable_docling: bool = Body(True),
    extract_tables: bool = Body(True),
    extract_images: bool = Body(True),
    ocr_enabled: bool = Body(False),
    output_format: str = Body("markdown")
):
    """파일 재벡터화 (기존 벡터 데이터 삭제 후 새로 벡터화)"""
    try:
        _clog.info(f"파일 재벡터화 요청: file_id={file_id}")
        
        # 파일 정보 확인
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        # VectorService 인스턴스 가져오기
        file_service = get_file_service_instance()
        if not file_service.vector_service:
            raise HTTPException(status_code=500, detail="VectorService를 사용할 수 없습니다.")
        
        # 기존 벡터 데이터 삭제
        try:
            await file_service.vector_service.delete_document_vectors(file_id)
            _clog.info(f"기존 벡터 데이터 삭제 완료: {file_id}")
        except Exception as e:
            _clog.warning(f"기존 벡터 데이터 삭제 중 오류 (계속 진행): {str(e)}")
        
        # Docling 옵션 설정
        docling_options = None
        if enable_docling:
            docling_options = DoclingOptions(
                output_format=output_format,
                extract_tables=extract_tables,
                extract_images=extract_images,
                ocr_enabled=ocr_enabled
            )
        
        # 벡터화 메타데이터 준비
        vector_metadata = {
            "filename": file_info.filename,
            "category_id": file_info.category_id,
            "category_name": file_info.category_name,
            "upload_time": file_info.upload_time.isoformat() if file_info.upload_time else None,
            "file_size": file_info.file_size
        }
        
        # 통합 벡터화 파이프라인 실행
        result = await file_service.vector_service.vectorize_with_docling_pipeline(
            file_path=file_info.file_path,
            file_id=file_id,
            metadata=vector_metadata,
            enable_docling=enable_docling,
            docling_options=docling_options
        )
        
        if result["success"]:
            # 파일 상태 업데이트
            await file_service.update_file_vectorization_status(
                file_id=file_id,
                vectorized=True,
                error_message=None,
                chunk_count=result["chunks_count"]
            )
            
            _ulog.info(
                "파일 재벡터화 완료",
                extra={
                    "event": "file_revectorization_completed",
                    "file_id": file_id,
                    "chunks_count": result["chunks_count"],
                    "processing_method": result.get("processing_method"),
                    "processing_time": result.get("processing_time")
                }
            )
            
            return {
                "success": True,
                "message": f"'{file_info.filename}' 파일이 성공적으로 재벡터화되었습니다.",
                "chunks_count": result["chunks_count"],
                "processing_method": result.get("processing_method", "unknown"),
                "processing_time": result.get("processing_time", 0),
                "docling_used": enable_docling
            }
        else:
            # 벡터화 실패 상태 업데이트
            await file_service.update_file_vectorization_status(
                file_id=file_id,
                vectorized=False,
                error_message=result.get("error", "재벡터화 실패"),
                chunk_count=0
            )
            
            _clog.error(f"파일 재벡터화 실패: {result.get('error')}")
            
            return {
                "success": False,
                "message": f"'{file_info.filename}' 파일 재벡터화에 실패했습니다.",
                "error": result.get("error", "알 수 없는 오류"),
                "docling_used": enable_docling
            }
            
    except HTTPException:
        raise
    except Exception as e:
        _clog.exception(f"파일 재벡터화 중 예기치 못한 오류: {str(e)}")
        raise HTTPException(status_code=500, detail=f"재벡터화 중 오류가 발생했습니다: {str(e)}")

@router.get("/{file_id}/download")
async def download_file(file_id: str):
    """파일 다운로드"""
    try:
        file_path = await get_file_service_instance().get_file_path(file_id)
        if not file_path:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        from fastapi.responses import FileResponse
        return FileResponse(
            path=file_path,
            filename=file_path.split('/')[-1],
            media_type='application/pdf'
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/{file_id}/view")
async def view_file(file_id: str):
    """파일 보기 - 다양한 파일 타입 지원"""
    try:
        # 파일 정보 조회
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        file_path = file_info.file_path
        if not file_path or not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        # 파일 확장자에 따른 MIME 타입 결정
        filename_lower = file_info.filename.lower()
        
        if filename_lower.endswith('.pdf'):
            media_type = 'application/pdf'
        elif filename_lower.endswith(('.jpg', '.jpeg')):
            media_type = 'image/jpeg'
        elif filename_lower.endswith('.png'):
            media_type = 'image/png'
        elif filename_lower.endswith('.gif'):
            media_type = 'image/gif'
        elif filename_lower.endswith('.bmp'):
            media_type = 'image/bmp'
        elif filename_lower.endswith(('.txt', '.md')):
            media_type = 'text/plain'
        elif filename_lower.endswith('.html'):
            media_type = 'text/html'
        elif filename_lower.endswith(('.doc', '.docx')):
            media_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        else:
            media_type = 'application/octet-stream'
        
        from fastapi.responses import FileResponse
        return FileResponse(
            path=file_path,
            media_type=media_type,
            filename=file_info.filename,
            headers={
                "Content-Disposition": f"inline; filename={file_info.filename}",
                "Cache-Control": "no-cache"
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/{file_id}/preview")
async def preview_file(file_id: str):
    """파일 미리보기 (브라우저 인라인 표시용)"""
    try:
        # 파일 정보 조회
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        file_path = file_info.file_path
        if not file_path or not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        # 파일 확장자에 따른 MIME 타입 결정
        filename_lower = file_info.filename.lower()
        
        if filename_lower.endswith('.pdf'):
            media_type = 'application/pdf'
        elif filename_lower.endswith(('.jpg', '.jpeg')):
            media_type = 'image/jpeg'
        elif filename_lower.endswith('.png'):
            media_type = 'image/png'
        elif filename_lower.endswith('.gif'):
            media_type = 'image/gif'
        elif filename_lower.endswith('.bmp'):
            media_type = 'image/bmp'
        elif filename_lower.endswith(('.txt', '.md')):
            media_type = 'text/plain; charset=utf-8'
        elif filename_lower.endswith('.html'):
            media_type = 'text/html; charset=utf-8'
        elif filename_lower.endswith('.csv'):
            media_type = 'text/csv; charset=utf-8'
        elif filename_lower.endswith('.docx'):
            media_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif filename_lower.endswith('.xlsx'):
            media_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif filename_lower.endswith('.pptx'):
            media_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        else:
            media_type = 'application/octet-stream'
        
        from fastapi.responses import FileResponse
        return FileResponse(
            path=file_path,
            media_type=media_type,
            headers={
                "Content-Disposition": "inline",
                "Cache-Control": "public, max-age=3600",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Methods": "GET",
                "Access-Control-Allow-Headers": "*"
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/{file_id}/content")
async def get_file_content(file_id: str):
    """텍스트 파일 내용 조회 (전처리 에디터용)"""
    try:
        # 파일 정보 조회
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        file_path = file_info.file_path
        if not file_path or not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        filename_lower = file_info.filename.lower()
        
        # 텍스트 파일인 경우 내용 읽기
        if filename_lower.endswith(('.txt', '.md', '.html', '.json', '.xml', '.csv')):
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                return {
                    "success": True,
                    "file_type": "text",
                    "content": content,
                    "filename": file_info.filename,
                    "file_size": file_info.file_size
                }
            except UnicodeDecodeError:
                # UTF-8로 읽기 실패시 다른 인코딩 시도
                try:
                    with open(file_path, 'r', encoding='cp949') as f:
                        content = f.read()
                    return {
                        "success": True,
                        "file_type": "text",
                        "content": content,
                        "filename": file_info.filename,
                        "file_size": file_info.file_size,
                        "encoding": "cp949"
                    }
                except:
                    return {
                        "success": False,
                        "error": "텍스트 파일을 읽을 수 없습니다. 인코딩 문제일 수 있습니다."
                    }
        # DOCX 파일인 경우 - 텍스트 추출
        elif filename_lower.endswith('.docx'):
            try:
                from docx import Document
                
                doc = Document(file_path)
                text_content = ""
                
                # 문단별로 텍스트 추출
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text_content += paragraph.text + "\n"
                
                # 표 내용도 추출
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text_content += " | ".join(row_text) + "\n"
                
                return {
                    "success": True,
                    "file_type": "docx",
                    "content": text_content.strip(),
                    "message": "DOCX에서 텍스트를 추출했습니다.",
                    "view_url": f"/api/v1/files/{file_id}/preview",
                    "filename": file_info.filename,
                    "file_size": file_info.file_size
                }
                
            except ImportError:
                return {
                    "success": False,
                    "error": "DOCX 텍스트 추출을 위해 python-docx 패키지가 필요합니다."
                }
            except Exception as e:
                return {
                    "success": False,
                    "error": f"DOCX 텍스트 추출 중 오류 발생: {str(e)}"
                }
        
        # XLSX 파일인 경우 - 구조화된 데이터 추출
        elif filename_lower.endswith(('.xlsx', '.xls')):
            try:
                from openpyxl import load_workbook
                import pandas as pd
                
                # openpyxl로 Excel 파일 읽기 (data_only=True로 수식 대신 값 가져오기)
                workbook = load_workbook(file_path, read_only=True, data_only=True)
                
                # 구조화된 데이터 저장
                sheets_data = {}
                text_content = ""  # 기존 텍스트 형태도 유지
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    
                    # 시트 데이터를 2차원 배열로 추출
                    sheet_data = []
                    max_row = 0
                    max_col = 0
                    
                    # 실제 데이터가 있는 영역 확인
                    for row in sheet.iter_rows():
                        row_data = []
                        has_data = False
                        for cell in row:
                            cell_value = cell.value
                            if cell_value is not None:
                                # 수식이 여전히 나오는 경우 빈 값으로 처리하거나 오류 메시지 표시
                                if isinstance(cell_value, str) and cell_value.startswith('='):
                                    # 수식이 계산되지 않은 경우 빈 값으로 처리
                                    row_data.append("[수식]")
                                else:
                                    row_data.append(str(cell_value))
                                has_data = True
                                max_col = max(max_col, len(row_data))
                            else:
                                row_data.append("")
                        
                        if has_data or sheet_data:  # 첫 데이터 이후 빈 행도 포함
                            sheet_data.append(row_data)
                            max_row = len(sheet_data)
                    
                    # 모든 행의 길이를 max_col로 맞춤
                    for row in sheet_data:
                        while len(row) < max_col:
                            row.append("")
                    
                    sheets_data[sheet_name] = {
                        "data": sheet_data,
                        "rows": max_row,
                        "cols": max_col
                    }
                    
                    # 텍스트 형태도 생성 (기존 청킹용)
                    text_content += f"\n=== {sheet_name} 시트 ===\n"
                    for row_data in sheet_data:
                        if any(cell.strip() for cell in row_data if cell):
                            text_content += " | ".join(row_data) + "\n"
                
                workbook.close()
                
                return {
                    "success": True,
                    "file_type": "xlsx",
                    "content": text_content.strip(),  # 텍스트 형태 (청킹용)
                    "sheets_data": sheets_data,       # 구조화된 데이터 (뷰어용)
                    "message": f"XLSX에서 데이터를 추출했습니다. ({len(sheets_data)}개 시트)",
                    "view_url": f"/api/v1/files/{file_id}/preview",
                    "filename": file_info.filename,
                    "file_size": file_info.file_size
                }
                
            except ImportError:
                return {
                    "success": False,
                    "error": "XLSX 데이터 추출을 위해 openpyxl 패키지가 필요합니다."
                }
            except Exception as e:
                return {
                    "success": False,
                    "error": f"XLSX 데이터 추출 중 오류 발생: {str(e)}"
                }
        
        # PPTX 파일인 경우 - 텍스트 추출
        elif filename_lower.endswith('.pptx'):
            try:
                from pptx import Presentation
                
                prs = Presentation(file_path)
                text_content = ""
                slide_count = len(prs.slides)
                
                for i, slide in enumerate(prs.slides, 1):
                    text_content += f"\n=== 슬라이드 {i} ===\n"
                    
                    # 슬라이드의 모든 텍스트 추출
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_content += shape.text + "\n"
                
                return {
                    "success": True,
                    "file_type": "pptx",
                    "content": text_content.strip(),
                    "message": f"PPTX에서 텍스트를 추출했습니다. ({slide_count}슬라이드)",
                    "view_url": f"/api/v1/files/{file_id}/preview",
                    "filename": file_info.filename,
                    "file_size": file_info.file_size,
                    "slide_count": slide_count
                }
                
            except ImportError:
                return {
                    "success": False,
                    "error": "PPTX 텍스트 추출을 위해 python-pptx 패키지가 필요합니다."
                }
            except Exception as e:
                return {
                    "success": False,
                    "error": f"PPTX 텍스트 추출 중 오류 발생: {str(e)}"
                }
        
        # PDF 파일인 경우 - 텍스트 추출 (AI 청킹 지원)
        elif filename_lower.endswith('.pdf'):
            try:
                # PDF에서 텍스트 추출 (PyMuPDF 사용)
                import fitz  # PyMuPDF
                
                doc = fitz.open(file_path)
                text_content = ""
                page_count = len(doc)
                
                for page_num in range(page_count):
                    page = doc[page_num]
                    page_text = page.get_text()
                    text_content += f"\n--- 페이지 {page_num + 1} ---\n{page_text}"
                
                doc.close()
                
                return {
                    "success": True,
                    "file_type": "pdf",
                    "content": text_content.strip(),
                    "message": f"PDF에서 텍스트를 추출했습니다. ({page_count}페이지)",
                    "view_url": f"/api/v1/files/{file_id}/view",
                    "filename": file_info.filename,
                    "file_size": file_info.file_size,
                    "page_count": page_count
                }
                
            except ImportError:
                return {
                    "success": False,
                    "error": "PDF 텍스트 추출을 위해 PyMuPDF 패키지가 필요합니다. 'pip install PyMuPDF'로 설치하세요."
                }
            except Exception as e:
                return {
                    "success": False,
                    "error": f"PDF 텍스트 추출 중 오류 발생: {str(e)}"
                }
        # 이미지 파일인 경우
        elif filename_lower.endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp')):
            return {
                "success": True,
                "file_type": "image",
                "message": "이미지 파일입니다.",
                "view_url": f"/api/v1/files/{file_id}/view",
                "filename": file_info.filename,
                "file_size": file_info.file_size
            }
        else:
            return {
                "success": False,
                "error": f"지원되지 않는 파일 형식입니다: {file_info.filename}"
            }
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/{file_id}/processing-status")
async def get_processing_status(file_id: str):
    """파일의 처리 상태 조회 (전처리 + 벡터화)"""
    try:
        file_info = await get_file_service_instance().get_file_info(file_id)
        if not file_info:
            raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
        
        return {
            "file_id": file_id,
            "filename": file_info.filename,
            "status": file_info.status,
            "upload_time": file_info.upload_time,
            "preprocessing_started_at": file_info.preprocessing_started_at,
            "preprocessing_completed_at": file_info.preprocessing_completed_at,
            "vectorization_started_at": file_info.vectorization_started_at,
            "vectorization_completed_at": file_info.vectorization_completed_at,
            "preprocessing_method": file_info.preprocessing_method,
            "chunk_count": file_info.chunk_count,
            "error_message": file_info.error_message,
            "vectorized": file_info.vectorized  # 하위 호환성
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/{file_id}/vectorization/retry")
async def retry_vectorization(file_id: str):
    """벡터화 재시도"""
    try:
        success = await get_file_service_instance().retry_vectorization(file_id)
        if not success:
            raise HTTPException(status_code=500, detail="벡터화 재시도에 실패했습니다.")
        return {"message": "벡터화 재시도가 시작되었습니다."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/search/")
async def search_documents(query: str, top_k: int = 5, category_ids: str = None):
    """문서 검색"""
    try:
        # 카테고리 ID 파싱
        category_list = None
        if category_ids:
            category_list = category_ids.split(',')
        
        # VectorService를 통한 검색 (지연 로딩)
        file_service = get_file_service_instance()
        if not file_service.vector_service:
            raise HTTPException(status_code=500, detail="벡터 검색 서비스를 초기화할 수 없습니다.")
        
        results = await file_service.vector_service.search_similar_chunks(
            query, top_k, category_list
        )
        
        return {
            "query": query,
            "results": results,
            "total_results": len(results)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/processing/status/")
async def get_processing_overview():
    """전체 파일 처리 상태 조회 (대시보드용)"""
    try:
        all_files = await get_file_service_instance().list_files()
        
        # 파일 상태별 통계
        status_stats = {
            "total_files": len(all_files),
            "uploaded": 0,
            "preprocessing": 0,
            "preprocessed": 0,
            "vectorizing": 0,
            "completed": 0,
            "failed": 0
        }
        
        # 메타데이터 기반으로 상태 집계
        files_summary = []
        
        for file_info in all_files:
            status = file_info.status.value if hasattr(file_info.status, 'value') else str(file_info.status)
            if status in status_stats:
                status_stats[status] += 1
            
            # 파일별 요약 정보 (메타데이터만 사용)
            files_summary.append({
                "file_id": file_info.file_id,
                "filename": file_info.filename,
                "status": status,
                "vectorized": file_info.vectorized,
                "category_name": file_info.category_name,
                "upload_time": file_info.upload_time,
                "file_size": file_info.file_size,
                "preprocessing_method": file_info.preprocessing_method,
                "chunk_count": file_info.chunk_count,
                "error_message": file_info.error_message
            })
        
        # ChromaDB 상태는 현재 초기화된 상태만 확인 (초기화 시도 없이)
        chromadb_status = "unknown"
        chromadb_message = "ChromaDB 상태를 확인하려면 별도 조회가 필요합니다."
        
        # 이미 초기화된 VectorService가 있다면 상태만 확인
        file_service = get_file_service_instance()
        if hasattr(file_service, '_vector_service') and file_service._vector_service:
            try:
                chroma_status = await file_service._vector_service.get_status()
                chromadb_status = chroma_status.get("connected", False)
                chromadb_message = "ChromaDB 연결됨" if chromadb_status else "ChromaDB 연결 안됨"
            except Exception:
                chromadb_status = False
                chromadb_message = "ChromaDB 상태 조회 중 오류 발생"
        else:
            chromadb_status = "not_loaded"
            chromadb_message = "ChromaDB가 로드되지 않았습니다. 벡터화 작업 시 자동으로 로드됩니다."
        
        return {
            "status_stats": status_stats,
            "files": files_summary,
            "chromadb_status": chromadb_status,
            "chromadb_message": chromadb_message,
            "total_completed_files": status_stats["completed"],
            "message": f"총 {len(all_files)}개 파일 중 {status_stats['completed']}개가 완전히 처리되었습니다."
        }
        
    except Exception as e:
        print(f"파일 처리 상태 조회 중 오류: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/vectorization/execute")
async def execute_vectorization(
    file_ids: Optional[List[str]] = Body(None),
    category_id: Optional[str] = Body(None)
):
    """선택된 파일들의 벡터화를 실행합니다."""
    try:
        # 설정 파일에서 기본 벡터화 Flow ID 읽기
        config_file = os.path.join(settings.BASE_DIR, "langflow", "config.json")
        vectorization_flow_id = None
        
        if os.path.exists(config_file):
            try:
                with open(config_file, 'r', encoding='utf-8') as f:
                    config_data = json.load(f)
                vectorization_flow_id = config_data.get("default_vectorization_flow_id")
            except Exception as e:
                print(f"설정 파일 읽기 오류: {str(e)}")
        
        if not vectorization_flow_id:
            raise HTTPException(
                status_code=400, 
                detail="LangFlow Flow ID가 설정되지 않았습니다. 관리자 설정에서 Flow ID를 설정해주세요."
            )
        
        # 대상 파일 목록 조회
        target_files = []
        if file_ids:
            # 특정 파일들만 벡터화
            for file_id in file_ids:
                file_info = await get_file_service_instance().get_file_info(file_id)
                if file_info and file_info.status in ["pending_vectorization", "vectorization_failed"]:
                    target_files.append(file_info)
        elif category_id:
            # 특정 카테고리의 모든 파일 벡터화 (최근 업로드된 것만)
            from datetime import datetime, timedelta
            recent_cutoff = datetime.now() - timedelta(hours=1)
            files = await get_file_service_instance().list_files(category_id)
            target_files = [f for f in files if f.status in ["uploaded", "pending_vectorization", "vectorization_failed"] and f.upload_time >= recent_cutoff]
        else:
            # 모든 대기 중인 파일 벡터화 (최근 업로드된 것만)
            from datetime import datetime, timedelta
            recent_cutoff = datetime.now() - timedelta(hours=1)
            all_files = await get_file_service_instance().list_files()
            target_files = [f for f in all_files if f.status in ["uploaded", "pending_vectorization", "vectorization_failed"] and f.upload_time >= recent_cutoff]
        
        if not target_files:
            _clog.debug("벡터화 대상 없음")
            return {"message": "벡터화할 파일이 없습니다."}
        
        # 백그라운드에서 벡터화 실행
        import asyncio
        for file_info in target_files:
            asyncio.create_task(get_file_service_instance().start_vectorization(file_info.file_id))
        _ulog.info(
            "벡터화 배치 시작",
            extra={"event": "vectorization_batch_started", "count": len(target_files)},
        )
        
        return {
            "message": f"{len(target_files)}개 파일의 벡터화가 시작되었습니다.",
            "target_files": [f.filename for f in target_files]
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/chromadb/status/")
async def get_chromadb_status():
    """ChromaDB 연결 및 상태 조회 (파일 API 통합 엔드포인트)"""
    try:
        # VectorService 직접 사용
        from ..services.vector_service import VectorService
        vector_service = VectorService()
        
        # ChromaDB 상태 조회
        status = await vector_service.get_status()
        return status
        
    except Exception as e:
        _clog.error(f"ChromaDB 상태 조회 실패: {e}")
        return {
            "connected": False,
            "total_vectors": 0,
            "collection_count": 0,
            "collections": [],
            "error": str(e)
        }

@router.get("/langflow/status/")
async def get_langflow_status():
    """LangFlow 등록 현황 조회"""
    try:
        # LangFlow 서비스를 직접 초기화하여 사용
        from ..services.langflow_service import LangflowService
        langflow_service = LangflowService()
        flows = await langflow_service.get_flows()
        
        # 설정 파일에서 기본 Flow ID 읽기
        config_file = os.path.join(settings.BASE_DIR, "langflow", "config.json")
        default_flow_id = None
        default_search_flow_id = None
        
        if os.path.exists(config_file):
            try:
                with open(config_file, 'r', encoding='utf-8') as f:
                    config_data = json.load(f)
                default_flow_id = config_data.get("default_vectorization_flow_id")
                default_search_flow_id = config_data.get("default_search_flow_id")
            except Exception as e:
                print(f"설정 파일 읽기 오류: {str(e)}")
        
        return {
            "default_vectorization_flow_id": default_flow_id,
            "default_search_flow_id": default_search_flow_id,
            "total_flows": len(flows) if flows else 0,
            "flows": flows or [],
            "langflow_configured": bool(default_flow_id or default_search_flow_id)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/langflow/flows/{flow_id}")
async def get_langflow_flow_details(flow_id: str):
    """특정 LangFlow Flow의 상세 정보 조회"""
    try:
        from ..services.langflow_service import LangflowService
        langflow_service = LangflowService()
        flow_details = await langflow_service.get_flow_details(flow_id)
        
        if not flow_details:
            raise HTTPException(status_code=404, detail="Flow를 찾을 수 없습니다.")
        
        return flow_details
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.put("/langflow/flows/{flow_id}/toggle")
async def toggle_flow_status(flow_id: str):
    """Flow 활성/비활성 상태를 토글합니다."""
    try:
        from ..services.langflow_service import LangflowService
        langflow_service = LangflowService()
        success = await langflow_service.toggle_flow_status(flow_id)
        
        if not success:
            raise HTTPException(status_code=404, detail="Flow를 찾을 수 없습니다.")
        
        return {"message": "Flow 상태가 변경되었습니다."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.put("/langflow/flows/{flow_id}/set-default")
async def set_default_vectorization_flow(flow_id: str):
    """Flow를 기본 벡터화 Flow로 설정합니다."""
    try:
        from ..services.langflow_service import LangflowService
        langflow_service = LangflowService()
        success = await langflow_service.set_default_vectorization_flow(flow_id)
        
        if not success:
            raise HTTPException(status_code=404, detail="Flow를 찾을 수 없습니다.")
        _ulog.info("기본 벡터화 Flow 설정", extra={"event": "langflow_set_default", "flow_id": flow_id})
        return {"message": "기본 벡터화 Flow가 설정되었습니다."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.put("/langflow/flows/{flow_id}/set-search")
async def set_search_flow(flow_id: str):
    """검색 Flow 설정"""
    try:
        success = await get_file_service_instance().set_search_flow(flow_id)
        if not success:
            raise HTTPException(status_code=404, detail="Flow를 찾을 수 없습니다.")
        _ulog.info("검색 Flow 설정", extra={"event": "langflow_set_search", "flow_id": flow_id})
        return {"message": f"검색 Flow가 {flow_id}로 설정되었습니다."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/langflow/flows/{flow_id}")
async def delete_flow(flow_id: str):
    """Flow 삭제"""
    try:
        success = await get_file_service_instance().delete_flow(flow_id)
        if not success:
            raise HTTPException(status_code=404, detail="Flow를 찾을 수 없습니다.")
        _ulog.info("Flow 삭제", extra={"event": "langflow_deleted", "flow_id": flow_id})
        return {"message": f"Flow {flow_id}가 성공적으로 삭제되었습니다."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



@router.post("/maintenance/diagnose-and-fix")
async def diagnose_and_fix_database():
    """ChromaDB 데이터베이스 진단 및 정상화를 순차적으로 실행합니다."""
    try:
        file_service = get_file_service_instance()
        vector_service = file_service.vector_service
        
        results = {
            "steps": [],
            "summary": {
                "total_issues_found": 0,
                "total_issues_fixed": 0,
                "status": "completed"
            }
        }
        
        # 1단계: 고아 메타데이터 검색 및 정리
        try:
            orphaned_metadata_count = await file_service.cleanup_orphaned_metadata()
            results["steps"].append({
                "step": "orphaned_metadata_cleanup",
                "name": "고아 메타데이터 정리",
                "status": "completed",
                "issues_found": orphaned_metadata_count,
                "issues_fixed": orphaned_metadata_count,
                "message": f"{orphaned_metadata_count}개의 고아 메타데이터를 정리했습니다."
            })
            results["summary"]["total_issues_found"] += orphaned_metadata_count
            results["summary"]["total_issues_fixed"] += orphaned_metadata_count
        except Exception as e:
            results["steps"].append({
                "step": "orphaned_metadata_cleanup",
                "name": "고아 메타데이터 정리",
                "status": "failed",
                "error": str(e),
                "message": f"고아 메타데이터 정리 중 오류: {str(e)}"
            })
        
        # 2단계: 고아 벡터 검색
        try:
            orphaned_vectors_info = await vector_service.find_orphaned_vectors()
            orphaned_vectors_count = orphaned_vectors_info.get('orphaned_count', 0)
            results["steps"].append({
                "step": "orphaned_vectors_detection",
                "name": "고아 벡터 검색",
                "status": "completed",
                "issues_found": orphaned_vectors_count,
                "message": f"{orphaned_vectors_count}개의 고아 벡터를 발견했습니다."
            })
            results["summary"]["total_issues_found"] += orphaned_vectors_count
        except Exception as e:
            results["steps"].append({
                "step": "orphaned_vectors_detection",
                "name": "고아 벡터 검색",
                "status": "failed",
                "error": str(e),
                "message": f"고아 벡터 검색 중 오류: {str(e)}"
            })
        
        # 3단계: 고아 벡터 정리
        try:
            cleanup_result = await vector_service.cleanup_orphaned_vectors()
            cleaned_count = cleanup_result.get('cleaned_count', 0)
            results["steps"].append({
                "step": "orphaned_vectors_cleanup",
                "name": "고아 벡터 정리",
                "status": "completed",
                "issues_fixed": cleaned_count,
                "message": f"{cleaned_count}개의 고아 벡터를 정리했습니다."
            })
            results["summary"]["total_issues_fixed"] += cleaned_count
        except Exception as e:
            results["steps"].append({
                "step": "orphaned_vectors_cleanup",
                "name": "고아 벡터 정리",
                "status": "failed",
                "error": str(e),
                "message": f"고아 벡터 정리 중 오류: {str(e)}"
            })
        
        # 4단계: 벡터화 상태 동기화
        try:
            sync_results = await file_service.sync_vectorization_status()
            status_corrected = sync_results.get('status_corrected', 0)
            results["steps"].append({
                "step": "vectorization_status_sync",
                "name": "벡터화 상태 동기화",
                "status": "completed",
                "issues_fixed": status_corrected,
                "message": f"{status_corrected}개 파일의 벡터화 상태를 동기화했습니다."
            })
            results["summary"]["total_issues_fixed"] += status_corrected
        except Exception as e:
            results["steps"].append({
                "step": "vectorization_status_sync",
                "name": "벡터화 상태 동기화",
                "status": "failed",
                "error": str(e),
                "message": f"벡터화 상태 동기화 중 오류: {str(e)}"
            })
        
        # 5단계: ChromaDB 상태 확인
        try:
            chromadb_status = await vector_service.get_status()
            results["steps"].append({
                "step": "chromadb_status_check",
                "name": "ChromaDB 상태 확인",
                "status": "completed",
                "chromadb_status": chromadb_status,
                "message": f"ChromaDB 상태: {'연결됨' if chromadb_status.get('connected') else '연결 안됨'}"
            })
        except Exception as e:
            results["steps"].append({
                "step": "chromadb_status_check",
                "name": "ChromaDB 상태 확인",
                "status": "failed",
                "error": str(e),
                "message": f"ChromaDB 상태 확인 중 오류: {str(e)}"
            })
        
        # 최종 요약 메시지 생성
        if results["summary"]["total_issues_found"] == 0:
            results["summary"]["message"] = "모든 진단이 완료되었습니다. 문제가 발견되지 않았습니다."
        else:
            results["summary"]["message"] = f"진단 및 정상화가 완료되었습니다. 총 {results['summary']['total_issues_found']}개 문제 중 {results['summary']['total_issues_fixed']}개를 해결했습니다."
        
        return results
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"진단 및 정상화 중 오류가 발생했습니다: {str(e)}")




@router.post("/chromadb/initialize")
async def initialize_chromadb():
    """ChromaDB를 수동으로 초기화합니다."""
    try:
        file_service = get_file_service_instance()
        vector_service = file_service.vector_service
        
        if not vector_service:
            return {
                "success": False,
                "error": "VectorService를 초기화할 수 없습니다.",
                "message": "VectorService 인스턴스 생성에 실패했습니다."
            }
        
        # ChromaDB 수동 초기화 시도
        result = await vector_service.initialize_chromadb_manually()
        
        if result["success"]:
            # 초기화 후 상태 확인
            status = await vector_service.get_status()
            result["status"] = status
            
        return result
        
    except Exception as e:
        error_msg = str(e)
        print(f"ChromaDB 초기화 API 오류: {error_msg}")
        return {
            "success": False,
            "error": error_msg,
            "message": f"초기화 중 API 오류가 발생했습니다: {error_msg}"
        } 